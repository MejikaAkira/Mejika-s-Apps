print("!!! EXECUTING app.py TOP LEVEL !!!")

# --- 標準ライブラリ ---
import asyncio
import io
import json
import os
import re
import sys
import gc
import csv
from collections import deque
from datetime import datetime

# --- サードパーティ ---
from dotenv import load_dotenv
from openai import AsyncOpenAI
import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_COLOR_TYPE
try:
    from pptx.enum.text import MSO_TEXT_DIRECTION
    HAS_TEXT_DIR = True
except ImportError:
    HAS_TEXT_DIR = False
from pptx.util import Inches, Pt

from quart import Quart, jsonify, render_template, request, send_file

# --- デバッグ出力 ---
print(f"DEBUG pptx version: {Presentation.__module__}")
print(f"DEBUG HAS_TEXT_DIR: {HAS_TEXT_DIR}")

# -------- 設定とグローバル変数 --------
# グローバル設定を保持する変数（キャッシュとして機能）
_config_cache = None

# 処理状態を追跡するグローバル変数
processing_status = {
    "is_processing": False,
    "current_file": "",
    "total_slides": 0,
    "current_slide": 0,
    "total_texts": 0,
    "translated_texts": 0,
    "stage": "idle",  # idle, loading, extracting, translating, summarizing, saving
    "message": "",
    "progress": 0.0,  # 0.0 - 100.0
}

# 翻訳キャッシュ
translation_cache = {}

# デフォルト設定
DEFAULT_CONFIG = {
    "api": {
        "base_url": "https://api.siemens.com/llm/v1",
        "model": "llama-3.1-8b-instruct"
    },
    "prompts": {},  # プロンプト詳細はconfig.jsonで管理
    # 翻訳から除外するテキストパターン
    "exclude_patterns": [
        r"Unrestricted \| © Siemens \d{4} \| Siemens Digital Industries Software",
        r"© Siemens \d{4} \| Siemens Digital Industries Software",
        r"Unrestricted",
        r"Page \d+",
        r"\d{1,2}/\d{1,2}/\d{2,4}",  # 日付形式
        r"\d{1,2}\.\d{1,2}\.\d{2,4}"  # 日付形式（ドット区切り）
    ]
}

# -------- 設定読み込み --------
def load_config(force_reload=False):
    """設定ファイルを読み込む"""
    global _config_cache
    
    # キャッシュがあり、強制再読み込みでなければキャッシュを返す
    if _config_cache is not None and not force_reload:
        return _config_cache
    
    # 設定ファイルのパスを取得
    config_dir = os.path.dirname(os.path.abspath(__file__))
    config_path = os.path.join(config_dir, 'config.json')
    
    # 設定ファイルの読み込み
    config_data = DEFAULT_CONFIG.copy()
    try:
        if os.path.exists(config_path):
            with open(config_path, 'r', encoding='utf-8') as f:
                loaded_config = json.load(f)
                # デフォルト設定に読み込んだ設定を上書き
                for key, value in loaded_config.items():
                    if key in config_data:
                        if isinstance(config_data[key], dict) and isinstance(value, dict):
                            config_data[key].update(value)
                        else:
                            config_data[key] = value
            print(f"設定ファイルを読み込みました: {config_path}")
        else:
            print(f"設定ファイルが見つかりません: {config_path}")
            print("デフォルト設定を使用します")
            # 設定ファイルが存在しない場合は作成
            try:
                os.makedirs(os.path.dirname(config_path), exist_ok=True)
                with open(config_path, 'w', encoding='utf-8') as f:
                    json.dump(DEFAULT_CONFIG, f, ensure_ascii=False, indent=2)
                print(f"デフォルト設定ファイルを作成しました: {config_path}")
            except Exception as e:
                print(f"設定ファイル作成エラー: {e}")
    except Exception as e:
        print(f"設定ファイル読み込みエラー: {e}")
        print("デフォルト設定を使用します")
    
    # 略語リスト（手動定義）
    abbreviations = ['NVH', 'E-motor', 'CAE', 'FEM', 'BEM', 'CFD', 'BEV', 'PHEV', 'ICE']
    config_data["abbreviations"] = abbreviations
    
    # キャッシュに保存
    _config_cache = config_data
    
    return config_data

# --- technical_terms.csv の読み込み ---
def load_technical_terms():
    """technical_terms.csvを読み込み、翻訳しない用語のリストを返す"""
    terms_dict = {}
    try:
        terms_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'technical_terms.csv')
        if os.path.exists(terms_file):
            with open(terms_file, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    if len(row) >= 1:  # 少なくとも1列は必要
                        term = row[0].strip()
                        replacement = row[1].strip() if len(row) >= 2 and row[1].strip() else term
                        terms_dict[term] = replacement
            add_log(f"technical_terms.csvを読み込みました: {len(terms_dict)}件の用語")
        else:
            add_log("technical_terms.csvが見つかりません。空のリストを使用します。")
    except Exception as e:
        add_log(f"technical_terms.csvの読み込みエラー: {e}")
    
    return terms_dict

# --- テキストクリーニング関数 ---
def should_exclude_text(text, shape_top=None):
    """翻訳から除外すべきテキストかどうかを判定する"""
    if not text or not text.strip():
        return True
    
    config_data = load_config()
    exclude_patterns = config_data.get("exclude_patterns", [])
    
    # 除外パターンに一致するかチェック
    for pattern in exclude_patterns:
        if re.search(pattern, text):
            return True
    
    # フッターっぽいテキストを除外（短くて下部にあるもの）
    if shape_top and shape_top > 5000000 and len(text) < 50:  # 経験的な閾値
        if any(keyword in text.lower() for keyword in ["siemens", "copyright", "©", "page", "unrestricted"]):
            return True
    
    return False

def fix_duplicate_abbreviations(text, abbreviations):
    """テキスト内の重複する略語を修正する"""
    if not text or not abbreviations:
        return text
    
    # 略語を長い順にソート（部分一致の問題を避けるため）
    sorted_abbrs = sorted(abbreviations, key=len, reverse=True)
    
    # 各略語について処理
    for abbr in sorted_abbrs:
        # 略語が2回以上出現する場合
        count = text.count(abbr)
        if count >= 2:
            # 最初の出現位置を取得
            first_pos = text.find(abbr)
            if first_pos >= 0:
                # 最初の出現以降のテキストを取得
                rest_text = text[first_pos + len(abbr):]
                
                # 残りのテキストから略語を検索
                for _ in range(count - 1):
                    next_pos = rest_text.find(abbr)
                    if next_pos >= 0:
                        # 前後のコンテキストを確認
                        before = text[:first_pos + len(abbr) + next_pos].strip()
                        
                        # 直前に同じ略語を含む単語がある場合、または括弧内の場合
                        if (before.endswith(abbr) or 
                            re.search(f"{re.escape(abbr)}[^a-zA-Z0-9]{{0,5}}$", before) or
                            (rest_text[next_pos-1:next_pos] in "（(" if next_pos > 0 else False)):
                            
                            # 重複している略語を削除
                            rest_text = rest_text[:next_pos] + rest_text[next_pos + len(abbr):]
                        else:
                            # 重複ではない場合は次の位置から検索
                            rest_text = rest_text[next_pos + len(abbr):]
                
                # 修正したテキストを再構成
                text = text[:first_pos + len(abbr)] + rest_text
    
    # 「NVH要素(NVH)」のようなパターンを検出
    for abbr in sorted_abbrs:
        pattern = f"([^a-zA-Z0-9])({re.escape(abbr)})[^a-zA-Z0-9]*\\([^\\)]*{re.escape(abbr)}[^\\)]*\\)"
        text = re.sub(pattern, r"\1\2", text)
    
    return text

def remove_translation_notes(text):
    """翻訳後の説明文を削除する"""
    if not text:
        return text
    
    # 1. 括弧内の説明文を削除するパターン
    patterns = [
        r'\(原文はそのまま訳し[^)]*\)',
        r'\(原文を直訳し[^)]*\)',
        r'\(略語や専門用語[^)]*\)',
        r'（原文はそのまま訳し[^）]*）',
        r'（原文を直訳し[^）]*）',
        r'（略語や専門用語[^）]*）',
        r'\(注[^)]*\)',
        r'（注[^）]*）',
    ]
    
    for pattern in patterns:
        text = re.sub(pattern, '', text)
    
    # 2. 文末に付く説明文を削除
    text = re.sub(r'(?:、|。)?[\s]*原文はそのまま訳し.*$', '', text)
    text = re.sub(r'(?:、|。)?[\s]*原文を直訳し.*$', '', text)
    
    # 3. 文頭に付く説明文を削除
    text = re.sub(r'^原文はそのまま訳し.*?(?=\S)', '', text)
    text = re.sub(r'^原文を直訳し.*?(?=\S)', '', text)
    text = re.sub(r'^はじめに(?:原文|翻訳)[^。]*', 'はじめに', text)
    
    # 4. 「〜したので、」などの接続表現を含む説明文を削除
    text = re.sub(r'原文はそのまま訳したので、.*?(?=\S)', '', text)
    text = re.sub(r'直訳したので、.*?(?=\S)', '', text)
    
    # 5. 文の途中に含まれる説明文を削除
    text = re.sub(r'(?<=\S)[\s]*原文はそのまま訳し[^。]*', '', text)
    text = re.sub(r'(?<=\S)[\s]*直訳し[^。]*', '', text)
    
    # 6. 翻訳プロセスに関する指示文を削除
    text = re.sub(r'(?:はじめに|ここでは)?(?:翻訳|訳文)(?:プロセス|処理)[^。]*(?:含め|追加)[^。]*', '', text)
    text = re.sub(r'(?:翻訳|訳文)(?:に|では)[^。]*(?:含め|追加)[^。]*', '', text)
    
    # 7. 「注:」「注意:」などを削除
    text = re.sub(r'^(?:注|注意|備考|補足)[:：].*?(?=\S)', '', text)
    text = re.sub(r'\((?:注|注意|備考|補足)[^)]*\)', '', text)
    text = re.sub(r'（(?:注|注意|備考|補足)[^）]*）', '', text)
    
    return text.strip()

def clean_instruction_text(text):
    """翻訳指示に関するテキストを削除する特別な関数"""
    if not text:
        return text
    
    # 「Agenda」の場合は特別処理
    if re.match(r'^Agenda[\s:：]?$', text.strip(), re.IGNORECASE):
        return 'Agenda'
    
    # 「はじめに」だけの場合は保持
    if text.strip() == 'はじめに':
        return 'はじめに'
    
    # 指示文のパターンを定義
    instruction_patterns = [
        # 翻訳プロセスに関する説明
        r'(?:はじめに|ここでは)?(?:翻訳|訳文)(?:プロセス|処理)[^。]*(?:含め|追加|しないで)[^。]*',
        r'(?:翻訳|訳文)(?:に|では)[^。]*(?:含め|追加|しないで)[^。]*',
        r'(?:原文|テキスト)[^。]*(?:そのまま|直訳)[^。]*(?:改行|説明|注釈)[^。]*(?:追加|含め)[^。]*',
        r'(?:説明|注釈)[^。]*(?:追加|含め)(?:ないで|しないで)[^。]*',
        r'はじめに(?:原文|翻訳)[^。]*',
        
        # より広範囲のパターン
        r'翻訳[^。]*(?:指示|命令|ガイドライン)[^。]*',
        r'(?:原文|テキスト)[^。]*(?:翻訳|訳)[^。]*(?:指示|命令|ガイドライン)[^。]*',
        r'(?:以下|下記)[^。]*(?:翻訳|訳)[^。]*(?:指示|命令|ガイドライン)[^。]*',
        r'(?:翻訳|訳)[^。]*(?:際|時)[^。]*(?:注意|留意)[^。]*',
        
        # 「原文をそのまま訳してください」などのパターン
        r'原文を(?:そのまま|直接)[^。]*(?:訳|翻訳)[^。]*(?:ください|下さい)',
        r'(?:訳|翻訳)[^。]*(?:際|時)[^。]*(?:原文|テキスト)[^。]*(?:そのまま|忠実)[^。]*',
        
        # 「説明や注釈を追加しないでください」などのパターン
        r'(?:説明|注釈|補足)[^。]*(?:追加|付け加え)(?:ないで|しないで)[^。]*(?:ください|下さい)',
        r'(?:説明|注釈|補足)[^。]*(?:不要|必要ない|省略)[^。]*',
        
        # 「略語はそのまま」などのパターン
        r'略語[^。]*(?:そのまま|変更しないで)[^。]*',
        r'(?:専門用語|技術用語)[^。]*(?:そのまま|説明なし)[^。]*',
        
        # Agendaスライドに残っている特定のパターン
        r'はじめに(?:翻訳|訳文)[^。]*関する[^。]*(?:絶対に)?(?:含め|追加)[^。]*(?:ないで|しないで)[^。]*',
        r'NVH要素の(?:翻訳|説明)[^。]*',
        r'(?:翻訳|訳文)(?:に|では)[^。]*NVH[^。]*',
        r'NVH(?:に|は)[^。]*(?:翻訳|説明)[^。]*',
        
        # このテキストは...系の説明
        r'(?:このテキスト|この文章|これ)は[^。]*(?:説明|表現|表す|意味|示す)[^。]*',
    ]
    
    # 各パターンを適用
    for pattern in instruction_patterns:
        text = re.sub(pattern, '', text)
    
    # 複数行のテキストの場合、各行ごとに処理
    if '\n' in text:
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # 各行が指示文っぽい場合は削除
            if (re.search(r'(?:翻訳|訳文|原文|説明|注釈)', line) and 
                len(line) > 10 and 
                not re.match(r'^(?:Agenda|はじめに|目次|概要)$', line.strip())):
                continue
            cleaned_lines.append(line)
        text = '\n'.join(cleaned_lines)
    
    # 「Agenda」だけの行は保持
    if re.match(r'^Agenda$', text.strip()):
        return 'Agenda'
    
    # 空になってしまった場合の処理
    if not text.strip() and 'agenda' in text.lower():
        return 'Agenda'
    
    return text.strip()

def clean_agenda_slide(text):
    """Agendaスライドの内容を特別にクリーニングする関数"""
    if not text:
        return text
    
    # 「Agenda」という単語だけを残す場合
    if re.match(r'^Agenda[\s:：]?$', text.strip(), re.IGNORECASE):
        return 'Agenda'
    
    # Agendaスライドの内容を行ごとに処理
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # 明らかに指示文と思われる行を削除
        if (re.search(r'(?:翻訳|訳文|原文|説明|注釈)(?:に|では|を|は|が|の)', line) or
            re.search(r'(?:含め|追加|しないで)(?:ください|下さい)', line)):
            continue
        
        # 短い項目（実際のアジェンダ項目と思われるもの）は保持
        if len(line.strip()) < 30 or not re.search(r'(?:翻訳|訳文|原文|説明|注釈)', line):
            cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines).strip()

def final_cleanup_check(text, is_agenda=False):
    """最終チェックとして明らかな指示文を削除する"""
    if not text:
        return text
    
    # Agendaの場合は特別処理
    if is_agenda and text.strip().lower() == "agenda":
        return "Agenda"
    
    # 明らかな指示文パターン
    obvious_instruction_patterns = [
        r'(?:翻訳|訳文)(?:プロセス|処理|に関する)[^。]*(?:絶対に)?(?:含め|追加|しないで)[^。]*',
        r'(?:原文|テキスト)(?:は|を)(?:そのまま|直訳)[^。]*',
        r'(?:説明|注釈)(?:は|を)(?:追加|含め)(?:ないで|しないで)[^。]*',
        r'はじめに(?:翻訳|原文)[^。]*',
        r'NVH要素[^。]*(?:翻訳|説明)[^。]*',
        r'(?:このテキスト|この文章|これ)は[^。]*(?:説明|表現|表す|意味|示す)[^。]*',
    ]
    
    # 各パターンを適用
    for pattern in obvious_instruction_patterns:
        text = re.sub(pattern, '', text)
    
    # 複数行テキストの場合、各行を個別にチェック
    if '\n' in text:
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # 明らかに指示文と思われる行を削除
            if (re.search(r'(?:翻訳|訳文|原文)(?:プロセス|処理|に関する)', line) or
                re.search(r'(?:含め|追加|しないで)(?:ください|下さい)', line)):
                continue
            cleaned_lines.append(line)
        text = '\n'.join(cleaned_lines)
    
    return text.strip()

def clean_title(title_text, original_text=None, abbreviations=None):
    """タイトル専用のクリーニング処理"""
    if not title_text:
        return title_text
    
    # 「Agenda」は特殊処理
    if original_text and original_text.strip().lower() == "agenda":
        return "Agenda"
    if re.match(r'^(?:アジェンダ|議題|予定|項目).*', title_text, re.IGNORECASE):
        return "Agenda"
    
    # 「このテキストは、タイトルであるため」などのメタ説明を削除
    title_text = re.sub(r'(?:このテキストは|これは)[^。]*(?:タイトル|見出し)[^。]*(?:ため|ので)[^。]*', '', title_text)
    
    # 「E-motor」「NVH」などの略語の後に続く説明を削除
    config_data = load_config()
    if abbreviations is None:
        abbreviations = config_data.get("abbreviations", [])
    
    for abbr in abbreviations:
        if original_text and abbr in original_text:
            # 略語に続く説明を削除（より積極的に）
            title_text = re.sub(f"{re.escape(abbr)}[^、。]*", abbr, title_text)
    
    # 括弧内の説明を削除
    title_text = re.sub(r'[\(（][^()（）]*?[\)）]', '', title_text)
    
    # 冒頭の「これは」「この」などの余分な言葉を削除
    title_text = re.sub(r'^(これは|この|ここでは|本資料では|本スライドでは)\s*', '', title_text)
    
    # 末尾の「です」「である」などを削除
    title_text = re.sub(r'(です|である|となります|します).*$', '', title_text)
    
    # 余分な記号や空白を整理
    title_text = re.sub(r'\s+', ' ', title_text).strip()
    title_text = re.sub(r'[、，,.:：；;]$', '', title_text).strip()
    
    # 重複する略語を修正
    title_text = fix_duplicate_abbreviations(title_text, abbreviations)
    
    # 「〜という意味」などの説明文を削除
    title_text = re.sub(r'(?:という|とは|は)[^。]*(?:意味|表す|示す|略)[^。]*', '', title_text)
    
    # 文の途中で切れている場合は、最初の句点までを取得
    if '。' in title_text:
        title_text = title_text.split('。')[0]
    
    # 翻訳後の説明文を削除
    title_text = remove_translation_notes(title_text)
    
    # 翻訳指示に関するテキストを削除
    title_text = clean_instruction_text(title_text)
    
    # 「注」「注意」などの単語を削除
    title_text = re.sub(r'^(?:注|注意|備考|補足)[:：]?', '', title_text)
    title_text = re.sub(r'\((?:注|注意|備考|補足)[^)]*\)', '', title_text)
    title_text = re.sub(r'（(?:注|注意|備考|補足)[^）]*）', '', title_text)
    
    # 最終クリーンアップ
    title_text = final_cleanup_check(title_text)
    
    return title_text.strip()

def clean_bullet_point(text, original_text=None, abbreviations=None):
    """箇条書きテキストのクリーニング処理"""
    if not text:
        return text
    
    # 「このテキストは」などのメタ説明を削除
    text = re.sub(r'(?:このテキストは|これは)[^。]*(?:箇条書き|リスト)[^。]*(?:ため|ので)[^。]*', '', text)
    
    # 略語の後に続く説明を削除
    config_data = load_config()
    if abbreviations is None:
        abbreviations = config_data.get("abbreviations", [])
    
    for abbr in abbreviations:
        if original_text and abbr in original_text:
            # 略語に続く説明を削除
            text = re.sub(f"{re.escape(abbr)}[（\(][^）\)]*[）\)]", abbr, text)
            text = re.sub(f"{re.escape(abbr)}は[^、。]*", f"{abbr}", text)
    
    # 「注:」などの注釈を削除
    text = re.sub(r'(?:^|\s)[（\(]?(?:注|Note|備考|補足|説明)[:：][^）\)]*?[）\)]?', '', text)
    
    # 「〜を表す」「〜の略」などのパターンを削除
    text = re.sub(r'は[^。]*(?:を表す|の略|を意味する|を示す)[^。]*', '', text)
    
    # 括弧内の説明を削除（原文に括弧がない場合）
    if original_text and not bool(re.search(r'[\(（].*?[\)）]', original_text)):
        text = re.sub(r'[\(（][^()（）]*?[\)）]', '', text)
    
    # 重複する略語を修正
    text = fix_duplicate_abbreviations(text, abbreviations)
    
    # 「〜という意味」などの説明文を削除
    text = re.sub(r'(?:という|とは)[^。]*(?:意味|表す|示す|略)[^。]*', '', text)
    
    # 翻訳後の説明文を削除
    text = remove_translation_notes(text)
    
    # 翻訳指示に関するテキストを削除
    text = clean_instruction_text(text)
    
    # 最終クリーンアップ
    text = final_cleanup_check(text)
    
    return text.strip()

def simplify_technical_text(text):
    """技術文書をよりシンプルでわかりやすい表現に変換する"""
    # 複雑な表現を簡潔な表現に置き換え
    replacements = [
        # 回りくどい表現の簡略化
        (r'パワートレインのマスクング処理が行われていない(?:ため|ことにより)', '内燃機関の音を隠す処理がないため'),
        (r'より明確に聞こえる道路音や風音が発生します', '道路や風の音がより目立ちます'),
        (r'([^、。]+)することが可能(?:です|になります)', r'\1できます'),
        (r'([^、。]+)する必要があります', r'\1してください'),
        (r'([^、。]+)と考えられます', r'\1と考えられます'),
        (r'([^、。]+)であると言えます', r'\1です'),
        
        # 長い修飾を簡略化
        (r'([^、。]+)するための([^、。]+)な([^、。]+)', r'\1用の\2\3'),
        
        # 受動態を能動態に
        (r'([^、。]+)によって([^、。]+)されます', r'\1が\2します'),
        
        # 冗長な表現の簡略化
        (r'([^、。]+)の観点から見ると', r'\1では'),
        (r'([^、。]+)という(?:観点|点)で', r'\1で'),
        
        # 二重否定の解消
        (r'([^、。]+)ないわけではありません', r'\1ます'),
        
        # 「〜のため」の簡略化
        (r'([^、。]+)を目的として', r'\1のため'),
    ]
    
    for pattern, replacement in replacements:
        text = re.sub(pattern, replacement, text)
    
    return text

def handle_technical_terms(text):
    """技術用語の翻訳を適切に処理する"""
    # 一般的な技術用語の自然な翻訳を確保
    replacements = {
        "ノイズ、振動、ハーシュネス": "NVH（ノイズ・振動・ハーシュネス）",
        "ノイズ、振動とハーシュネス": "NVH（ノイズ・振動・ハーシュネス）",
        "ノイズ振動ハーシュネス": "NVH",
        "電気自動車": "BEV",
        "プラグインハイブリッド車": "PHEV",
        "内燃機関": "ICE"
    }
    
    for orig, replacement in replacements.items():
        if orig.lower() in text.lower():
            text = re.sub(re.escape(orig), replacement, text, flags=re.IGNORECASE)
    
    return text

def improve_naturalness(text):
    """翻訳文の自然さを向上させる（強化版）"""
    # 不自然な表現を修正
    replacements = [
        (r'([の])である([。])', r'\1\2'),  # 「〜のである。」→「〜の。」
        (r'([すまいる])ます([。])', r'\1\2'),  # 「〜します。」→「〜す。」など
        (r'([のだ])です([。])', r'\1\2'),  # 「〜のだです。」→「〜のだ。」
        (r'電気モーターの?NVH', 'E-モーターのNVH'),  # 一貫性のために
        (r'([0-9]+)パーセント', r'\1%'),  # 数字+パーセント → 数字+%
        # --- 追加: よくある不自然な日本語パターン ---
        (r'することができます', 'できます'),
        (r'することが可能です', 'できます'),
        (r'する必要があります', 'してください'),
        (r'である。', 'です。'),
        (r'である,', 'です,'),
        (r'である ', 'です '),
        (r'である$', 'です'),
        (r'\s+', ' '),  # 連続スペースを1つに
        (r'\n+', '\n'),  # 連続改行を1つに
        (r'\s+。', '。'),
        (r'\s+,', ','),
        (r'\s+、', '、'),
        (r'\s+：', '：'),
        (r'\s+:', ':'),
        (r'。+', '。'),  # 句点の連続を1つに
        (r'、+', '、'),  # 読点の連続を1つに
        (r'\s+$', ''),  # 文末スペース除去
        (r'^\s+', ''),  # 文頭スペース除去
    ]
    for pattern, replacement in replacements:
        text = re.sub(pattern, replacement, text)
    return text.strip()

# --- ユーティリティ関数 ---
from pptx.enum.dml import MSO_COLOR_TYPE

def get_original_font_color(run):
    if hasattr(run.font, 'color') and run.font.color:
        color = run.font.color
        # 明示的なRGB指定
        if hasattr(color, 'rgb') and color.rgb:
            return color.rgb
        # それ以外（Auto/None）はNone扱い
    return None

def get_shape_bg_rgb(shape):
    """シェイプの背景色を取得する関数（改良版）"""
    try:
        # シェイプ自体の背景色を確認
        if hasattr(shape, 'fill') and shape.fill:
            fill_type = getattr(shape.fill, 'type', None)
            # _NoFillタイプの場合はスキップ
            if fill_type is not None and fill_type != 0:  # 0 = MSO_FILL.NO_FILL
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color and hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                    return shape.fill.fore_color.rgb
        # スライドの背景色を確認
        if hasattr(shape, 'part') and hasattr(shape.part, 'slide') and hasattr(shape.part.slide, 'background'):
            bg = shape.part.slide.background
            if hasattr(bg, 'fill') and bg.fill:
                bg_fill_type = getattr(bg.fill, 'type', None)
                if bg_fill_type is not None and bg_fill_type != 0:  # 0 = MSO_FILL.NO_FILL
                    if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color and hasattr(bg.fill.fore_color, 'rgb') and bg.fill.fore_color.rgb:
                        return bg.fill.fore_color.rgb
    except Exception as e:
        add_log(f"背景色取得エラー: {e}")
    # 背景色が取得できない場合は白を仮定
    return RGBColor(255, 255, 255)

def is_bright_color(rgb):
    """色の明るさを判定する関数（改良版）"""
    if rgb is None:
        return True  # 取得できなければ白背景扱い
    if isinstance(rgb, int):
        r = (rgb >> 16) & 0xFF
        g = (rgb >> 8) & 0xFF
        b = rgb & 0xFF
    else:
        try:
            r = rgb.r if hasattr(rgb, 'r') else 255
            g = rgb.g if hasattr(rgb, 'g') else 255
            b = rgb.b if hasattr(rgb, 'b') else 255
        except:
            return True
    luminance = (0.299 * r + 0.587 * g + 0.114 * b)
    return luminance > 128

def set_text_color_if_needed(run, shape):
    """フォント色が設定されていない場合のみ、背景色に基づいてテキスト色を設定"""
    if hasattr(run.font, 'color') and run.font.color and hasattr(run.font.color, 'type'):
        # MSO_COLOR_TYPE.RGB (1) の場合は既に色が設定されている
        if getattr(run.font.color, 'type', None) == 1:
            return  # 既に色が設定されているので何もしない
    try:
        bg_rgb = get_shape_bg_rgb(shape)
        if is_bright_color(bg_rgb):
            run.font.color.rgb = RGBColor(0, 0, 0)
        else:
            run.font.color.rgb = RGBColor(255, 255, 255)
    except Exception as e:
        add_log(f"テキスト色設定エラー: {e}")
        try:
            run.font.color.rgb = RGBColor(0, 0, 0)
        except:
            pass

# -------- .env 読み込み & OpenAI 初期化 --------
load_dotenv(override=True)  # 環境変数を確実に上書き

try:
    # 設定を読み込む
    config_data = load_config()
    
    # SSL 検証を無効にしたカスタム HTTP クライアントを作成
    http_client = httpx.AsyncClient(
        verify=os.getenv("OPENAI_SSL_VERIFY", "true").lower() != "false",
        timeout=httpx.Timeout(300)  # タイムアウトを300秒に設定
    )

    # 環境変数からベース URL を取得（設定ファイルの値をデフォルトとして使用）
    base_url = os.environ.get("OPENAI_BASE_URL", config_data["api"]["base_url"])
    
    # 起動時に環境変数の状態を表示
    print(f"API Base URL: {base_url}")
    print(f"Model: {config_data['api']['model']}")
    
    client = AsyncOpenAI(
        api_key=os.environ.get("OPENAI_API_KEY"),
        base_url=base_url,
        http_client=http_client
    )
    print("API client initialized with SSL verification disabled.")
except Exception as e:
    client = None
    print(f"API initialization failed: {e}")

# -------- Quart アプリ --------
app = Quart(__name__)
app.config["PROVIDE_AUTOMATIC_OPTIONS"] = True  # CORS対策
app.config["MAX_CONTENT_LENGTH"] = 512 * 1024 * 1024  # 512MBに制限を増やす

# -------- ログ --------
LOG_BUFFER_SIZE = 200
log_buffer = deque(maxlen=LOG_BUFFER_SIZE)
def add_log(msg): 
    print(msg); log_buffer.append(msg); sys.stdout.flush()
    # 処理状態にメッセージを追加
    processing_status["message"] = msg

# -------- OpenAI ヘルパ --------
async def generate_summary(texts: list[dict], max_tokens=800) -> str:
    """スライドの内容からエグゼクティブサマリーを生成する"""
    if not client:
        return "Error: OpenAI client not initialized."
    
    processing_status["stage"] = "summarizing"
    add_log("サマリーを生成中...")
    
    # テキストが空の場合は早期リターン
    if not texts:
        add_log("サマリー生成用のテキストがありません")
        return "スライド内容が見つかりませんでした。"
    
    config_data = load_config()
    
    # スライド内容の全体像を把握するための情報を集約
    slide_titles = []
    slide_contents = {}
    
    for i, text_info in enumerate(texts):
        slide_id = text_info["id"]
        slide_text = text_info["text"]
        
        # 空のテキストはスキップ
        if not slide_text.strip():
            continue
        
        # スライドタイトルを抽出する試み（最初の短いテキストをタイトルと仮定）
        lines = slide_text.split('\n')
        potential_title = lines[0] if lines else ""
        if len(potential_title) < 100:  # タイトルらしき短いテキスト
            slide_titles.append(f"{slide_id}: {potential_title}")
        else:
            slide_titles.append(slide_id)
        
        # スライド内容を保存
        slide_contents[slide_id] = slide_text
    
    # スライド構造の概要を作成
    slide_structure = "\n".join(slide_titles)
    
    # サマリー生成用のスライド内容が十分にあるか確認
    if len(slide_contents) < 2:
        add_log("サマリー生成用のスライドが不足しています")
        return "スライドの内容が不十分なため、有意義なサマリーを生成できません。"
    
    # より良いサマリー生成のためのプロンプト
    improved_prompt = f"""
    あなたはプレゼンテーション資料を分析し、高品質なエグゼクティブサマリーを作成する専門家です。
    
    以下のプレゼンテーションを分析し、経営層向けの簡潔で価値のあるエグゼクティブサマリーを作成してください。
    
    【プレゼンテーション構造】
    {slide_structure}
    
    【スライド詳細内容】
    {json.dumps(slide_contents, ensure_ascii=False, indent=2)}
    
    エグゼクティブサマリーの作成ガイドライン:
    1. 全体で400-600字程度にまとめてください
    2. プレゼンテーション全体から得られる主要な知見、結論、提案を抽出してください
    3. 単なるスライド内容の羅列ではなく、情報を統合して価値ある洞察を提供してください
    4. 専門用語（NVH、E-motorなど）はそのまま使用してください
    5. 文章形式で作成し、箇条書きは使わないでください
    6. 「エグゼクティブサマリー」などのタイトルは含めないでください
    7. 段落を2-3つ程度に分けて読みやすくしてください
    
    エグゼクティブサマリー:
    """
    
    try:
        # 改善されたサマリーを生成
        rsp = await client.chat.completions.create(
            model=config_data["api"]["model"],
            messages=[
                {"role": "system", "content": "あなはプレゼンテーション資料を分析し、経営層向けの価値あるサマリーを作成する専門家です。"},
                {"role": "user", "content": improved_prompt}
            ],
            max_tokens=max_tokens,
            temperature=0.4
        )
        
        summary = rsp.choices[0].message.content.strip()
        
        # サマリーから冒頭の「エグゼクティブサマリー」などのタイトルを削除
        summary = re.sub(r'^[\s\n]*(?:エグゼクティブ)?(?:サマリー|概要)[:：]?[\s\n]*', '', summary)
        summary = re.sub(r'^[\s\n]*Executive Summary[:：]?[\s\n]*', '', summary, flags=re.IGNORECASE)
        
        add_log("サマリー生成完了")
        return summary
    except Exception as e:
        add_log(f"サマリー生成エラー: {e}")
        return f"サマリー生成中にエラーが発生しました: {str(e)}"

async def translate_text(txt: str) -> str:
    """テキストを翻訳する（キャッシュ機能付き）"""
    global translation_cache
    
    # 空のテキストはそのまま返す
    if not client or not txt.strip():
        return txt
    
    # キャッシュにあれば、それを返す
    if txt in translation_cache:
        # キャッシュヒットのカウントを更新
        processing_status["translated_texts"] += 1
        processing_status["progress"] = min(90, (processing_status["translated_texts"] / processing_status["total_texts"]) * 90)
        
        # 定期的に進捗をログに記録
        if processing_status["translated_texts"] % 20 == 0 or processing_status["translated_texts"] == processing_status["total_texts"]:
            add_log(f"翻訳進捗 (キャッシュ使用): {processing_status['translated_texts']}/{processing_status['total_texts']} ({processing_status['progress']:.1f}%)")
        
        return translation_cache[txt]
    
    config_data = load_config()
    translation_config = config_data["prompts"]["translation"]
    abbreviations = config_data.get("abbreviations", [])
    
    # 記号・箇条書きだけの場合はスキップ
    if re.fullmatch(r'^[\s\u2022\u25AA\u25CF\u25B6■●\-–—]+$', txt):
        return txt
    
    # 単語「Agenda」は翻訳せずにそのまま残す
    if txt.strip().lower() == "agenda":
        translation_cache[txt] = "Agenda"  # キャッシュに保存
        return "Agenda"
    
    # technical_terms.csvから翻訳しない用語を読み込む
    technical_terms = load_technical_terms()
    
    # 元のテキストに改行が含まれているかチェック
    has_newlines = '\n' in txt
    is_title = len(txt) < 50  # 50文字未満はタイトルと見なす
    
    # 翻訳しない用語をマークアップ
    marked_text = txt
    term_placeholders = {}
    
    # 製品名や専門用語を一時的にプレースホルダーに置き換え
    for i, (term, replacement) in enumerate(technical_terms.items()):
        if term.lower() in marked_text.lower():
            placeholder = f"__TERM_{i}__"
            # 大文字小文字を区別せずに置換
            pattern = re.compile(re.escape(term), re.IGNORECASE)
            marked_text = pattern.sub(placeholder, marked_text)
            term_placeholders[placeholder] = replacement
    
    # プロンプトの構築
    prompt = translation_config["user"]
    
    # タイトルの場合と本文の場合で異なる指示を追加
    if is_title:
        prompt += " これはタイトルなので、改行を入れずに翻訳してください。"
    elif has_newlines:
        prompt += " 元のテキストの改行パターンを尊重してください。"
    
    # 翻訳しない用語のリストがある場合は指示に追加
    if term_placeholders:
        terms_list = ", ".join([f'"{p}" → "{v}"' for p, v in term_placeholders.items()])
        prompt += f"\n\n以下の専門用語やプレースホルダーは翻訳せず、指定された形式をそのまま使用してください：\n{terms_list}"
    
    prompt += f"\n\n```\n{marked_text}\n```"
    
    try:
        system_message = translation_config["system"]
        
        if is_title:
            system_message += " タイトルの場合は改行を入れずに翻訳してください。"
        elif has_newlines:
            system_message += " 元のテキストの改行パターンを尊重してください。"
        
        rsp = await client.chat.completions.create(
            model=config_data["api"]["model"],
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": prompt},
            ],
            max_tokens=int(len(txt) * 1.5) + 20,
            temperature=translation_config.get("temperature", 0.1),
        )
        out = rsp.choices[0].message.content.strip()
        
        # 基本的な後処理
        # 1. コードブロックを除去
        if out.startswith("```") and out.endswith("```"):
            stripped = out[3:-3].strip()
            if "\n" in stripped:
                # 最初の行が言語タグの場合は除去
                _, _, stripped = stripped.partition("\n")
            out = stripped.strip()
        
        # 2. 引用符を除去
        if out.startswith('"') and out.endswith('"'):
            out = out[1:-1]
        
        # 3. 日本語の括弧を除去
        if (out.startswith("「") and out.endswith("」")) or (out.startswith("『") and out.endswith("』")):
            out = out[1:-1]
        
        # 4. プレースホルダーを元の用語に戻す
        for placeholder, term in term_placeholders.items():
            out = out.replace(placeholder, term)
        
        # 5. 余分な空白を削除
        out = re.sub(r'([^\s])\s+([^\s])', r'\1\2', out)
        
        # 6. タイトルの場合は改行を削除
        if is_title and '\n' in out:
            out = out.replace('\n', ' ')
        
        # 7. 略語の後に続く説明を削除
        if abbreviations:
            # 略語リストから正規表現パターンを作成
            abbr_pattern = '|'.join(map(re.escape, abbreviations))
            # 略語の後に続く括弧内の説明を削除
            out = re.sub(f'({abbr_pattern})[\s]*[（\(][^）\)]*[）\)]', r'\1', out)
            # 「〜は〜の略」などのパターンを削除
            for abbr in abbreviations:
                out = re.sub(f'{abbr}は[^、。]*(?:を表す|の略|を意味する|を示す)[^、。]*[、。]?', f'{abbr}', out)
        
        # 8. 「注:」や「Note:」で始まる注釈を削除
        out = re.sub(r'(?:^|\s)[（\(]?(?:注|Note|備考|補足|説明)[:：][^）\)]*?[）\)]?', '', out)
        
        # 9. 文字数に関する言及を削除
        out = re.sub(r'(原文|翻訳)[^。]*文字数[^。]*。?', '', out)
        out = re.sub(r'この(翻訳|訳文)[^。]*文字[^。]*。?', '', out)
        out = re.sub(r'文字数[^。]*維持[^。]*。?', '', out)
        out = re.sub(r'簡潔に訳[^。]*。?', '', out)
        
        # 10. 「は〜を表す」「は〜の略」などのパターンを削除
        out = re.sub(r'は[^。]*(?:を表す|の略|を意味する|を示す)[^。]*', '', out)
        
        # 11. 「このテキストは」などのメタ説明を削除
        out = re.sub(r'(?:このテキストは|これは)[^。]*(?:タイトル|見出し|箇条書き)[^。]*(?:ため|ので)[^。]*', '', out)
        
        # 12. 「〜という意味です」などの説明を削除
        out = re.sub(r'(?:という|とは)[^。]*(?:意味|表す|示す|略)[^。]*(?:です|である|します).*$', '', out)
        
        # 13. 重複する略語を修正
        out = fix_duplicate_abbreviations(out, abbreviations)
        
        # 14. 翻訳後の説明文を削除
        out = remove_translation_notes(out)
        
        # 15. 翻訳指示に関するテキストを削除
        out = clean_instruction_text(out)
        
        # 16. 「注」「(注」などを削除
        out = re.sub(r'^(?:注|注意|備考|補足)[:：]?', '', out)
        out = re.sub(r'\((?:注|注意|備考|補足)[^)]*\)', '', out)
        out = re.sub(r'（(?:注|注意|備考|補足)[^）]*）', '', out)
        
        # 17. Agendaスライドの場合は特別処理
        if "agenda" in txt.lower():
            out = clean_agenda_slide(out)
        
        # 18. 最終クリーンアップチェック
        is_agenda_slide = "agenda" in txt.lower() or "agenda" in out.lower()
        out = final_cleanup_check(out, is_agenda=is_agenda_slide)
        
        # 19. 技術用語の処理
        out = handle_technical_terms(out)
        
        # 20. 自然さの向上
        out = improve_naturalness(out)
        
        # 21. 技術文書の簡略化
        out = simplify_technical_text(out)
        
        # 翻訳結果をキャッシュに保存
        translation_cache[txt] = out.strip()
        
        # 翻訳済みテキスト数を更新
        processing_status["translated_texts"] += 1
        processing_status["progress"] = min(90, (processing_status["translated_texts"] / processing_status["total_texts"]) * 90)
        
        # 定期的に進捗をログに記録
        if processing_status["translated_texts"] % 20 == 0 or processing_status["translated_texts"] == processing_status["total_texts"]:
            add_log(f"翻訳進捗: {processing_status['translated_texts']}/{processing_status['total_texts']} ({processing_status['progress']:.1f}%)")
        
        return out.strip()
    except Exception as e:
        print(f"Translation error: {e}")
        return txt  # エラーが発生した場合は元のテキストを返す

# スライドを一度に少しずつ処理するための関数
async def process_slides_in_batches(prs, batch_size=5):
    """スライドをバッチ処理して、メモリ使用量を最適化する"""
    total_slides = len(prs.slides)
    processing_status["total_slides"] = total_slides
    
    all_text_runs = []
    all_texts_for_summary = []
    
    for batch_start in range(0, total_slides, batch_size):
        batch_end = min(batch_start + batch_size, total_slides)
        add_log(f"スライド {batch_start+1}-{batch_end} を処理中 (全{total_slides}スライド)")
        processing_status["current_slide"] = batch_end
        
        batch_text_runs = []
        batch_texts_for_summary = []
        
        # このバッチのスライドを処理
        for si in range(batch_start, batch_end):
            slide = prs.slides[si]
            slide_txt_collect = []
            
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                
                # テキストフレームの位置を取得
                shape_top = shape.top
                
                # directionを水平に統一
                if HAS_TEXT_DIR:
                    try:
                        shape.text_frame.direction = MSO_TEXT_DIRECTION.HORIZONTAL
                    except Exception:
                        pass
                
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        txt = run.text
                        
                        # 除外すべきテキストはスキップ
                        if should_exclude_text(txt, shape_top):
                            continue
                        
                        if txt.strip():
                            slide_txt_collect.append(txt)
                            batch_text_runs.append(run)
            
            if slide_txt_collect:
                batch_texts_for_summary.append({"id": f"Slide {si+1}", "text": " ".join(slide_txt_collect)})
        
        # 全体のリストに追加
        all_text_runs.extend(batch_text_runs)
        all_texts_for_summary.extend(batch_texts_for_summary)
        
        # メモリを解放するためにガベージコレクションを実行
        gc.collect()
    
    # 総テキスト数を設定
    processing_status["total_texts"] = len(all_text_runs)
    add_log(f"テキスト抽出完了。翻訳対象: {len(all_text_runs)}個のテキストボックス")
    
    return all_text_runs, all_texts_for_summary

# サマリースライドを作成する関数
def create_summary_slides(prs, summary_text):
    """サマリーテキストからスライドを作成する（オリジナルのスタイルを継承）"""
    add_log("サマリースライドを作成中...")
    
    # サマリーから冒頭の「エグゼクティブサマリー」などのタイトルを削除
    summary_text = re.sub(r'^[\s\n]*(?:エグゼクティブ)?(?:サマリー|概要)[:：]?[\s\n]*', '', summary_text)
    summary_text = re.sub(r'^[\s\n]*Executive Summary[:：]?[\s\n]*', '', summary_text, flags=re.IGNORECASE)
    
    # サマリーを段落に分割
    paragraphs = []
    current_paragraph = []
    for line in summary_text.split('\n'):
        line = line.strip()
        # 各行からも「エグゼクティブサマリー」を削除
        line = re.sub(r'^[\s\n]*(?:エグゼクティブ)?(?:サマリー|概要)[:：]?[\s\n]*', '', line)
        line = re.sub(r'^[\s\n]*Executive Summary[:：]?[\s\n]*', '', line, flags=re.IGNORECASE)
        
        if not line:  # 空行は段落の区切り
            if current_paragraph:
                paragraphs.append(' '.join(current_paragraph))
                current_paragraph = []
        else:
            current_paragraph.append(line)
    
    # 最後の段落を追加
    if current_paragraph:
        paragraphs.append(' '.join(current_paragraph))
    
    # 段落がない場合は元のテキストをそのまま使用
    if not paragraphs:
        paragraphs = [summary_text]
    
    # 1ページあたりの最大文字数
    MAX_CHARS_PER_PAGE = 800
    
    # サマリーを複数ページに分割
    current_page_content = ""
    page_contents = []
    
    for paragraph in paragraphs:
        # このパラグラフを追加すると最大文字数を超える場合は新しいページに
        if len(current_page_content) + len(paragraph) > MAX_CHARS_PER_PAGE and current_page_content:
            page_contents.append(current_page_content)
            current_page_content = paragraph
        else:
            if current_page_content:
                current_page_content += "\n\n" + paragraph
            else:
                current_page_content = paragraph
    
    # 最後のページを追加
    if current_page_content:
        page_contents.append(current_page_content)
    
    # ページがない場合は1ページ追加
    if not page_contents:
        page_contents = ["サマリー情報はありません。"]
    
    # 各ページを作成
    total_pages = len(page_contents)
    summary_slides = []
    
    # オリジナルのスライドからスタイルをコピー
    template_slide = None
    if len(prs.slides) > 0:
        template_slide = prs.slides[0]  # 最初のスライドをテンプレートとして使用
    
    for i, content in enumerate(page_contents):
        # スライドを作成（オリジナルと同じレイアウトを使用）
        layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
        summary_slide = prs.slides.add_slide(layout)
        summary_slides.append(summary_slide)
        
        # タイトル用テキストボックス
        left, top = Inches(0.5), Inches(0.5)
        width = Inches(9)
        box_t = summary_slide.shapes.add_textbox(left, top, width, Inches(0.75))
        p = box_t.text_frame.paragraphs[0]
        title_text = "エグゼクティブサマリー"
        if total_pages > 1:
            title_text += f" ({i+1}/{total_pages})"
        p.text = title_text
        
        # オリジナルのスタイルを適用（タイトル）
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        
        # オリジナルのスライドからフォント色をコピー
        if template_slide:
            try:
                for shape in template_slide.shapes:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                                title_run = shape.text_frame.paragraphs[0].runs[0]
                                if hasattr(title_run.font, 'color') and hasattr(title_run.font.color, 'rgb'):
                                    p.font.color.rgb = title_run.font.color.rgb
                                break
            except:
                pass
        
        # 本文用テキストボックス
        box_b = summary_slide.shapes.add_textbox(
            left, 
            top + Inches(1),  # タイトルの下に配置
            width, 
            Inches(6)  # 十分な高さを確保
        )
        
        tf = box_b.text_frame
        tf.word_wrap = True
        
        # 本文テキストを追加
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(12)
        p.space_after = Pt(12)  # 段落後の間隔
        
        # オリジナルのスライドから本文のフォント色をコピー
        if template_slide:
            try:
                for shape in template_slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        # タイトル以外のテキストを探す
                        for para in shape.text_frame.paragraphs:
                            if para.runs and len(para.text) > 10:  # 本文らしきテキスト
                                body_run = para.runs[0]
                                if hasattr(body_run.font, 'color') and hasattr(body_run.font.color, 'rgb'):
                                    p.font.color.rgb = body_run.font.color.rgb
                                break
            except:
                pass
    
    add_log(f"サマリースライドを {total_pages} ページ作成しました")
    return summary_slides

# -------- ルート --------
@app.route('/')
async def index():
    tpl = os.path.join('templates', 'index.html')
    if os.path.exists(tpl):
        return await render_template('index.html')
    return '<p>index.html が見つかりません。</p>', 404

@app.route('/logs')
async def logs():
    return jsonify({"logs": list(log_buffer)})

@app.route('/status')
async def status():
    """現在の処理状態を返す"""
    return jsonify(processing_status)

@app.route('/upload', methods=['POST'])
async def upload():
    global processing_status, translation_cache
    
    if not client:
        return "OpenAI 未初期化", 503
    
    # 既に処理中の場合はエラーを返す
    if processing_status["is_processing"]:
        return "別のファイルを処理中です。完了までお待ちください。", 409

    add_log("アップロードリクエストを受信しました")
    
    try:
        files = await request.files  # ※Quart では await 必須
        add_log(f"リクエストファイル: {files}")
        
        if 'file' not in files:
            add_log("ファイルがリクエストに含まれていません")
            return 'ファイルがありません', 400
        
        f = files['file']
        if f.filename == '':
            add_log("ファイル名が空です")
            return 'ファイル名が空です', 400
        
        if not f.filename.lower().endswith('.pptx'):
            add_log(f"無効なファイル形式: {f.filename}")
            return 'pptx をアップしてください', 400

        # 処理状態を初期化
        processing_status = {
            "is_processing": True,
            "current_file": f.filename,
            "total_slides": 0,
            "current_slide": 0,
            "total_texts": 0,
            "translated_texts": 0,
            "stage": "loading",
            "message": "ファイルを読み込み中...",
            "progress": 0.0
        }
        
        # 翻訳キャッシュをクリア
        translation_cache = {}

        fname = f.filename
        
        # ファイルデータを読み込む
        try:
            data = f.read()  # FileStorage.read は同期
            add_log(f"Received {fname}, {len(data)} bytes")
        except Exception as e:
            add_log(f"ファイル読み込みエラー: {str(e)}")
            processing_status["is_processing"] = False
            return f"ファイル読み込みエラー: {str(e)}", 500

        try:
            # --- オリジナルファイルをサーバーに保存 ---
            UPLOAD_DIR = 'uploads'
            os.makedirs(UPLOAD_DIR, exist_ok=True)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            original_path = os.path.join(UPLOAD_DIR, f"original_{timestamp}_{fname}")
            with open(original_path, "wb") as f_out:
                f_out.write(data)
            add_log(f"Original saved to {original_path}")

            # --- コピーを編集用として開く ---
            processing_status["stage"] = "loading"
            add_log("PowerPointファイルを開いています...")
            prs = Presentation(original_path)
            add_log(f"PowerPointを開きました。{len(prs.slides)}枚のスライドがあります。")

            # --- 元テキスト抽出（フッターなどを除外） ---
            processing_status["stage"] = "extracting"
            add_log("スライドからテキストを抽出中...")
            
            shapes_to_translate = []
            texts_for_summary = []
            for slide in prs.slides:
                slide_txt_collect = []
                for shape in slide.shapes:
                    if getattr(shape, 'has_text_frame', False):
                        shapes_to_translate.append(shape)
                        shape_text = "\n".join([para.text for para in shape.text_frame.paragraphs]).strip()
                        if shape_text:
                            slide_txt_collect.append(shape_text)
                if slide_txt_collect:
                    texts_for_summary.append({"id": f"Slide {prs.slides.index(slide)+1}", "text": " ".join(slide_txt_collect)})

            processing_status["total_texts"] = sum(len(shape.text_frame.paragraphs) for shape in shapes_to_translate)
            add_log(f"テキスト抽出完了。翻訳対象: {processing_status['total_texts']}個の段落")

            # --- 段落単位で翻訳 ---
            processing_status["stage"] = "translating"
            add_log(f"翻訳を開始します。対象段落数: {processing_status['total_texts']}")
            
            paragraphs_to_translate = []
            para_shape_refs = []  # (shape, para)のペア
            for shape in shapes_to_translate:
                for para in shape.text_frame.paragraphs:
                    if not hasattr(para, 'runs') or not para.runs:
                        continue
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        paragraphs_to_translate.append(text)
                        para_shape_refs.append((shape, para))

            concurrency = int(os.getenv("OPENAI_CONCURRENCY", "5"))
            sem = asyncio.Semaphore(concurrency)
            async def sem_tr(t): 
                async with sem: 
                    return await translate_text(t)
            
            trans_tasks = [asyncio.create_task(sem_tr(text)) for text in paragraphs_to_translate]
            translations = await asyncio.gather(*trans_tasks)

            add_log("翻訳が完了しました。テキストを適用中...")
            # --- シンプルなrun置換でスタイルを維持 ---
            for (shape, para), translated_text in zip(para_shape_refs, translations):
                if not para or not hasattr(para, 'runs') or not para.runs:
                    continue
                runs = para.runs
                runs[0].text = translated_text
                for run in runs[1:]:
                    run.text = ''

            # --- 保存 ---
            processing_status["stage"] = "saving"
            add_log("翻訳済みPowerPointを保存中...")
            
            # 結果を一時ファイルに保存
            result_path = os.path.join(UPLOAD_DIR, f"translated_{timestamp}_{fname}")
            prs.save(result_path)
            
            # 結果をメモリに読み込み
            buf = io.BytesIO()
            with open(result_path, "rb") as f_in:
                buf.write(f_in.read())
            buf.seek(0)
            
            processing_status["progress"] = 100.0
            processing_status["stage"] = "completed"
            add_log("処理が完了しました！")
            add_log(f"キャッシュ使用率: {len(translation_cache)}/{processing_status['total_texts']} ({len(translation_cache)/max(1, processing_status['total_texts'])*100:.1f}%)")

            # --- Quartバージョンに応じてファイル名指定用引数を切り替え ---
            import inspect
            send_file_params = inspect.signature(send_file).parameters
            kwargs_common = dict(
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
            )
            
            # 処理状態をリセット
            processing_status["is_processing"] = False
            
            if 'download_name' in send_file_params:
                # 新バージョン (>=0.19) 用
                return await send_file(buf, download_name=f"translated_{fname}", **kwargs_common)
            else:
                # 旧バージョン (<=0.18) 用
                return await send_file(buf, attachment_filename=f"translated_{fname}", **kwargs_common)
        
        except Exception as e:
            # エラーが発生した場合
            processing_status["stage"] = "error"
            processing_status["message"] = f"エラーが発生しました: {str(e)}"
            processing_status["is_processing"] = False
            add_log(f"エラー: {str(e)}")
            return f"処理中にエラーが発生しました: {str(e)}", 500
    
    except Exception as e:
        # アップロード処理自体でエラーが発生した場合
        processing_status["is_processing"] = False
        add_log(f"アップロードエラー: {str(e)}")
        return f"アップロード中にエラーが発生しました: {str(e)}", 500

# -------- サーバ起動 --------
if __name__ == '__main__':
    print("Starting Quart server on http://127.0.0.1:8000")
    from hypercorn.config import Config
    from hypercorn.asyncio import serve
    
    hypercorn_config = Config()
    hypercorn_config.bind = ["127.0.0.1:8000"]
    hypercorn_config.keep_alive_timeout = 3600  # 60分
    hypercorn_config.worker_class = "asyncio"
    hypercorn_config.request_timeout = 3600  # リクエストタイムアウト60分
    
    # メモリ制限を増やす試み（Windows では利用できないのでスキップ）
    try:
        import platform
        if platform.system() != "Windows":  # Windows 以外の場合のみ実行
            import resource
            resource.setrlimit(resource.RLIMIT_AS, (resource.RLIM_INFINITY, resource.RLIM_INFINITY))
            print("メモリ制限を無制限に設定しました")
    except Exception as e:
        print(f"メモリ制限の設定に失敗しました: {e}")
    
    asyncio.run(serve(app, hypercorn_config))

def load_terms_csv():
    terms = {'term': set(), 'no_translate': set(), 'remove': set(), 'replace': {}}
    terms_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'technical_terms.csv')
    with open(terms_file, encoding='utf-8') as f:
        reader = csv.DictReader(f)
        for row in reader:
            t = row['type'].strip()
            pat = row['pattern'].strip()
            rep = row['replacement'].strip() if 'replacement' in row else ''
            if t == 'term':
                terms['term'].add(pat)
            elif t == 'no_translate':
                terms['no_translate'].add(pat)
            elif t == 'remove':
                terms['remove'].add(pat)
            elif t == 'replace':
                terms['replace'][pat] = rep
    return terms

def protect_terms(text, terms):
    """Replace all term/no_translate words with unique placeholders."""
    placeholders = {}
    all_terms = list(terms.get('term', set()) | terms.get('no_translate', set()))
    # 長い語順で置換（部分一致防止）
    all_terms.sort(key=len, reverse=True)
    for i, term in enumerate(all_terms):
        if not term:
            continue
        placeholder = f"__TERM_{i}__"
        text = re.sub(re.escape(term), placeholder, text)
        placeholders[placeholder] = term
    return text, placeholders

def restore_terms(text, placeholders):
    for placeholder, term in placeholders.items():
        text = text.replace(placeholder, term)
    return text

def clean_text(text, terms=None):
    if terms is None:
        terms = load_terms_csv()
    # 1. term/no_translate語を一時保護
    text, placeholders = protect_terms(text, terms)
    # 2. 除去
    for pat in terms['remove']:
        text = re.sub(re.escape(pat), '', text)
    # 3. 置換
    for pat, rep in terms['replace'].items():
        text = re.sub(re.escape(pat), rep, text)
    # 4. term/no_translate語を復元
    text = restore_terms(text, placeholders)
    return text

# 既存の専門用語・クリーニング・置換・除去処理は全てclean_textで統一して呼び出すようにしてください。
# 例: translated_text = clean_text(translated_text)

def postprocess_terms(original, translated, terms):
    """
    For each term/no_translate word present in the original, ensure it is present in the translation.
    If missing, append it (or in the future, replace likely substitutions). This is a minimal, robust coexistence approach.
    """
    for term in terms:
        if term in original and term not in translated:
            # Append the term if missing (future: replace synonyms/variants)
            translated = translated.strip() + f"（{term}）"
    return translated

# --- upload関数の翻訳適用ループ内で、term/no_translate語が含まれる場合のみpostprocess_termsを適用してください ---
# 例:
#   target_terms = [t for t in (terms['term'] | terms['no_translate']) if t in original_text]
#   if target_terms:
#       translated_text = postprocess_terms(original_text, translated_text, target_terms)

def build_translation_prompt(base_prompt, chunk, terms):
    """
    If the chunk contains any term/no_translate words, append an instruction to the prompt
    to preserve them exactly as in the original. This is generic and works for any term list.
    """
    target_terms = [t for t in (terms['term'] | terms['no_translate']) if t in chunk]
    if target_terms:
        joined = ', '.join(target_terms)
        base_prompt += f'\n\n次の単語・フレーズは必ず原文のまま訳文に残してください: {joined}'
    return base_prompt

# --- translate_textや翻訳バッチ処理の直前で、build_translation_promptを使ってプロンプトを拡張してください ---
# 例:
#   prompt = build_translation_prompt(base_prompt, chunk, terms)
#   translated = await translate_text(chunk, prompt)