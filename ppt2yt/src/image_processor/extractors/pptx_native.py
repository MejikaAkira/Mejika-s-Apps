"""
python-pptxを使用したスライド抽出（互換性改善版）
"""
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional
from PIL import Image, ImageDraw, ImageFont
import io

# collections.abc 互換性対応
if sys.version_info >= (3, 10):
    import collections.abc
    collections.Iterable = collections.abc.Iterable
    collections.Mapping = collections.abc.Mapping
    collections.MutableMapping = collections.abc.MutableMapping
    collections.MutableSet = collections.abc.MutableSet
    collections.Sequence = collections.abc.Sequence

from .base import BaseSlideExtractor


class NativePptxExtractor(BaseSlideExtractor):
    """python-pptxを使用したスライド抽出"""
    
    def is_available(self) -> bool:
        """python-pptxが利用可能かチェック"""
        try:
            # 互換性チェックを含む
            from pptx import Presentation
            from PIL import Image
            
            # テスト用の空のプレゼンテーションを作成してみる
            try:
                prs = Presentation()
                return True
            except Exception as e:
                self.logger.warning(f"python-pptx互換性チェックエラー: {e}")
                return False
                
        except ImportError:
            return False
    
    def can_extract_video(self) -> bool:
        """動画抽出をサポートしていない"""
        return False
    
    def extract_single_slide(self, pptx_path: str, slide_number: int) -> Optional[Dict[str, Any]]:
        """単一スライドを画像として抽出"""
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            if slide_number <= len(prs.slides):
                slide = prs.slides[slide_number - 1]
                image = self._render_slide_to_image(slide)
                
                if image:
                    filename = f"slide_{slide_number:02d}.png"
                    output_path = self.output_dir / filename
                    image.save(output_path, "PNG")
                    
                    if output_path.exists():
                        return self._create_slide_info(slide_number, output_path)
                    else:
                        self.logger.warning(f"スライド {slide_number} の保存に失敗")
                else:
                    self.logger.warning(f"スライド {slide_number} のレンダリングに失敗")
            else:
                self.logger.error(f"スライド番号 {slide_number} は存在しません")
                
        except Exception as e:
            self.logger.error(f"単一スライド抽出でエラー: {e}")
            return None
    
    def extract_slides(self, pptx_path: str) -> List[Dict[str, Any]]:
        """すべてのスライドを画像として抽出"""
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            images_info = []
            
            for i, slide in enumerate(prs.slides, 1):
                try:
                    image = self._render_slide_to_image(slide)
                    
                    if image:
                        filename = f"slide_{i:02d}.png"
                        output_path = self.output_dir / filename
                        image.save(output_path, "PNG")
                        
                        if output_path.exists():
                            images_info.append(self._create_slide_info(i, output_path))
                        else:
                            self.logger.warning(f"スライド {i} の保存に失敗")
                    else:
                        self.logger.warning(f"スライド {i} のレンダリングに失敗")
                        
                except Exception as e:
                    self.logger.error(f"スライド {i} の処理でエラー: {e}")
                    continue
            
            return images_info
        except Exception as e:
            self.logger.error(f"スライド抽出でエラー: {e}")
            return []
    
    def extract_slide_as_video(self, pptx_path: str, slide_number: int, 
                              duration: int = 10) -> Optional[Path]:
        """スライドを動画として抽出（サポートしていない）"""
        self.logger.warning("python-pptx抽出器は動画抽出をサポートしていません")
        return None
    
    def _render_slide_to_image(self, slide) -> Optional[Image.Image]:
        """スライドを画像にレンダリング（改善版）"""
        try:
            # スライドサイズを安全に取得
            try:
                slide_width = slide.slide_width
                slide_height = slide.slide_height
            except AttributeError:
                # デフォルト値を使用
                slide_width = 9144000  # 10インチ
                slide_height = 6858000  # 7.5インチ
            
            # 設定からサイズを取得
            width = self.config.get("export_width", 1920)
            height = self.config.get("export_height", 1080)
            
            # 画像を作成
            image = Image.new('RGB', (width, height), 'white')
            draw = ImageDraw.Draw(image)
            
            # スライドの内容を描画
            self._draw_slide_content_safe(draw, slide, width, height, slide_width, slide_height)
            
            return image
        except Exception as e:
            self.logger.error(f"スライドレンダリングでエラー: {e}")
            return None
    
    def _draw_slide_content_safe(self, draw, slide, width: int, height: int, 
                                slide_width: int, slide_height: int):
        """スライドの内容を安全に描画"""
        try:
            # 背景を描画
            self._draw_background(draw, slide, width, height)
            
            # シェイプを分類
            text_shapes = []
            image_shapes = []
            other_shapes = []
            
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text_shapes.append(shape)
                    elif hasattr(shape, 'shape_type'):
                        shape_type = shape.shape_type
                        # MSO_SHAPE_TYPE の値を直接使用（互換性のため）
                        if shape_type == 13:  # PICTURE
                            image_shapes.append(shape)
                        else:
                            other_shapes.append(shape)
                except Exception as e:
                    self.logger.debug(f"シェイプ分類エラー: {e}")
            
            # 各種シェイプを描画
            for shape in other_shapes:
                self._draw_shape_safe(draw, shape, width, height, slide_width, slide_height)
            
            for shape in image_shapes:
                self._draw_image_shape(draw, shape, width, height, slide_width, slide_height)
            
            for shape in text_shapes:
                self._draw_text_shape_safe(draw, shape, width, height, slide_width, slide_height)
                
        except Exception as e:
            self.logger.error(f"スライド内容描画でエラー: {e}")
    
    def _draw_text_shape_safe(self, draw, shape, width: int, height: int,
                            slide_width: int, slide_height: int):
        """テキストシェイプを安全に描画"""
        try:
            if not hasattr(shape, 'text_frame'):
                return
            
            text = shape.text_frame.text
            if not text:
                return
            
            # 位置とサイズを安全に取得
            left = self._safe_convert_emu(shape.left, width, slide_width)
            top = self._safe_convert_emu(shape.top, height, slide_height)
            
            # フォントを取得
            try:
                # システムフォントを試行
                font = ImageFont.truetype("arial.ttf", 24)
            except:
                font = ImageFont.load_default()
            
            # テキストを描画
            draw.text((left, top), text, fill='black', font=font)
            
        except Exception as e:
            self.logger.debug(f"テキスト描画エラー: {e}")
    
    def _safe_convert_emu(self, emu_value, pixel_size: int, slide_size: int) -> int:
        """EMU値を安全にピクセルに変換"""
        try:
            if emu_value is None:
                return 0
            return int(emu_value * pixel_size / slide_size)
        except:
            return 0
    
    def _draw_background(self, draw, slide, width: int, height: int):
        """背景を描画"""
        try:
            # スライドの背景色を取得（可能な場合）
            fill_color = 'white'  # デフォルト
            
            # 背景を塗りつぶし
            draw.rectangle([0, 0, width, height], fill=fill_color)
        except Exception as e:
            self.logger.debug(f"背景描画エラー: {e}")
    
    def _draw_image_shape(self, draw, shape, width: int, height: int,
                         slide_width: int, slide_height: int):
        """画像シェイプを描画"""
        try:
            if hasattr(shape, 'image'):
                # 画像データを取得
                image_stream = io.BytesIO(shape.image.blob)
                img = Image.open(image_stream)
                
                # 位置とサイズを計算
                left = self._safe_convert_emu(shape.left, width, slide_width)
                top = self._safe_convert_emu(shape.top, height, slide_height)
                shape_width = self._safe_convert_emu(shape.width, width, slide_width)
                shape_height = self._safe_convert_emu(shape.height, height, slide_height)
                
                # 画像をリサイズ
                img = img.resize((shape_width, shape_height), Image.Resampling.LANCZOS)
                
                # 画像を描画（エラーを回避）
                try:
                    draw._image.paste(img, (left, top))
                except:
                    # フォールバック：枠を描画
                    draw.rectangle([left, top, left + shape_width, top + shape_height],
                                 outline='gray', width=2)
                    
        except Exception as e:
            self.logger.debug(f"画像描画エラー: {e}")
    
    def _draw_shape_safe(self, draw, shape, width: int, height: int,
                        slide_width: int, slide_height: int):
        """その他のシェイプを安全に描画"""
        try:
            # 位置とサイズを取得
            left = self._safe_convert_emu(shape.left, width, slide_width)
            top = self._safe_convert_emu(shape.top, height, slide_height)
            shape_width = self._safe_convert_emu(shape.width, width, slide_width)
            shape_height = self._safe_convert_emu(shape.height, height, slide_height)
            
            # 基本的な矩形を描画
            draw.rectangle([left, top, left + shape_width, top + shape_height],
                         outline='lightgray', width=1)
                         
        except Exception as e:
            self.logger.debug(f"シェイプ描画エラー: {e}") 