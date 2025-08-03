"""
スライド内のメディア（動画）検出
"""
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Dict, Any, Optional

from ...utils.logger import get_logger


class MediaDetector:
    """スライド内のメディア検出"""
    
    def __init__(self):
        self.logger = get_logger("MediaDetector")
        self.video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv']
    
    def detect_videos_in_slide(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """スライド内の動画を検出"""
        videos = []
        
        # 1. PPTXファイル構造から検出
        videos.extend(self._detect_from_pptx_structure(pptx_path, slide_number))
        
        # 2. python-pptxで検出
        videos.extend(self._detect_from_pptx_library(pptx_path, slide_number))
        
        # 重複を除去
        return self._remove_duplicates(videos)
    
    def _detect_from_pptx_structure(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """PPTXファイル構造から動画を検出"""
        videos = []
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                # スライドのリレーションシップを確認
                slide_rels_path = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
                
                if slide_rels_path in pptx_zip.namelist():
                    rels_content = pptx_zip.read(slide_rels_path)
                    rels_tree = ET.fromstring(rels_content)
                    
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    
                    for rel in rels_tree.findall('.//r:Relationship', ns):
                        target = rel.get('Target')
                        if target and '../media/' in target:
                            media_filename = target.split('/')[-1]
                            
                            if any(media_filename.lower().endswith(ext) for ext in self.video_extensions):
                                videos.append({
                                    "media_file": media_filename,
                                    "rel_id": rel.get('Id'),
                                    "source": "pptx_structure"
                                })
                                
        except Exception as e:
            self.logger.error(f"PPTXファイル構造の検査でエラー: {e}")
            
        return videos
    
    def _detect_from_pptx_library(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """python-pptxライブラリで動画を検出"""
        videos = []
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(pptx_path)
            if slide_number <= len(prs.slides):
                slide = prs.slides[slide_number - 1]
                
                for shape in slide.shapes:
                    # メディアタイプのシェイプを確認
                    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                        videos.append({
                            "shape_name": shape.name,
                            "shape_type": "media",
                            "source": "pptx_library"
                        })
                        
        except Exception as e:
            self.logger.error(f"python-pptxでの動画検出でエラー: {e}")
            
        return videos
    
    def _remove_duplicates(self, videos: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """重複を除去"""
        unique_videos = []
        seen = set()
        
        for video in videos:
            key = video.get('media_file') or video.get('shape_name')
            if key and key not in seen:
                seen.add(key)
                unique_videos.append(video)
                
        return unique_videos
    
    def extract_embedded_video(self, pptx_path: str, slide_number: int, 
                             video_info: Dict[str, Any], output_dir: Path) -> Optional[Path]:
        """埋め込み動画を抽出"""
        try:
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                media_file = video_info.get('media_file')
                
                if media_file:
                    media_path = f'ppt/media/{media_file}'
                    
                    if media_path in pptx_zip.namelist():
                        video_data = pptx_zip.read(media_path)
                        
                        # 出力ファイル名
                        extension = Path(media_file).suffix
                        video_filename = f"slide_{slide_number:02d}_embedded{extension}"
                        video_path = output_dir / video_filename
                        
                        # ファイルとして保存
                        with open(video_path, 'wb') as f:
                            f.write(video_data)
                            
                        self.logger.info(f"埋め込み動画を抽出: {video_path}")
                        return video_path
                        
        except Exception as e:
            self.logger.error(f"埋め込み動画の抽出でエラー: {e}")
            
        return None 