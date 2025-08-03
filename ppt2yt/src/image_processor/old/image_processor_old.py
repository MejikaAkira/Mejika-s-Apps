"""
画像処理モジュール
PPTXファイルから画像を抽出し、動画用に処理
"""
import os
import json
import io
import subprocess
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

from ..utils.config import config
from ..utils.logger import get_logger


class ImageProcessor:
    """PPTXファイルから画像を抽出・処理するクラス"""
    
    def __init__(self):
        """ImageProcessorの初期化"""
        self.logger = get_logger("ImageProcessor")
        self.output_resolution = config.get("video.resolution", "1920x1080")
        self.image_format = config.get("image.format", "png")
        self.image_quality = config.get("image.quality", 95)
        
        # 動画処理設定
        self.extract_embedded_videos = config.get("video_processing.extract_embedded_videos", True)
        self.video_format = config.get("video_processing.video_format", "mp4")
        self.gif_fps = config.get("video_processing.gif_fps", 10)
        self.gif_quality = config.get("video_processing.gif_quality", 85)
        self.max_video_duration = config.get("video_processing.max_video_duration", 30)
        self.video_scale = config.get("video_processing.video_scale", "640x360")
        self.use_unified_processing = config.get("video_processing.use_unified_processing", True)
        self.save_slides_as_video = config.get("video_processing.save_slides_as_video", True)
        self.extract_embedded_videos_properly = config.get("video_processing.extract_embedded_videos_properly", True)
        
        # 解像度をパース
        self.width, self.height = self._parse_resolution(self.output_resolution)
        self.video_width, self.video_height = self._parse_resolution(self.video_scale)
        
        # 出力ディレクトリ
        self.output_dir = config.get_path("paths.output.images")
        if self.output_dir:
            self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def _parse_resolution(self, resolution: str) -> Tuple[int, int]:
        """
        解像度文字列をパース
        
        Args:
            resolution: 解像度文字列（例: "1920x1080"）
            
        Returns:
            (width, height)のタプル
        """
        try:
            width, height = map(int, resolution.split('x'))
            return width, height
        except ValueError:
            self.logger.warning(f"解像度のパースに失敗: {resolution}。デフォルトの1920x1080を使用")
            return 1920, 1080
    
    def extract_slides_as_images(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        PPTXファイルからスライド全体を画像として抽出（3段階フォールバック）
        
        Args:
            pptx_path: PPTXファイルのパス
            
        Returns:
            画像情報のリスト（全スライドを画像化）
        """
        self.logger.info(f"PPTXファイルからスライド全体を画像化中: {pptx_path}")
        
        # 1. PowerPoint COM Export（最高品質）
        if self._can_use_win32com():
            self.logger.info("PowerPoint COM Export を使用して画像を抽出します")
            result = self._extract_with_powerpoint(pptx_path)
            if result:
                return result
        
        # 2. PowerPoint SaveAs（高品質）
        if self._can_use_win32com():
            self.logger.info("PowerPoint SaveAs を使用して画像を抽出します")
            result = self._extract_with_powerpoint_saveas(pptx_path)
            if result:
                return result
        
        # 3. LibreOffice（クロスプラットフォーム）
        if self._can_use_libreoffice():
            self.logger.info("LibreOffice を使用して画像を抽出します")
            result = self._extract_with_libreoffice(pptx_path)
            if result:
                return result
        
        # 4. フォールバック（標準的な方法）
        self.logger.info("標準的な方法で画像を抽出します")
        return self._extract_slides_with_pil(pptx_path)
    
    def extract_slides_with_video_check(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        スライドごとに動画の有無を確認し、適切な形式で保存
        
        Args:
            pptx_path: PPTXファイルのパス
            
        Returns:
            画像・動画情報のリスト
        """
        self.logger.info(f"スライドごとの動画確認処理を開始: {pptx_path}")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)
            images_info = []
            
            for slide_number in range(1, total_slides + 1):
                self.logger.info(f"=== スライド {slide_number}/{total_slides} の処理開始 ===")
                
                # スライド内の動画を確認
                videos = self.extract_embedded_videos_from_slide(pptx_path, slide_number)
                
                if videos:
                    # 動画がある場合の処理
                    if self.save_slides_as_video:
                        # スライド全体をMP4として保存
                        video_path = self._save_slide_as_video(pptx_path, slide_number)
                        
                        if video_path:
                            images_info.append({
                                "slide_number": slide_number,
                                "image_path": str(video_path),
                                "filename": video_path.name,
                                "description": f"スライド {slide_number} （動画）",
                                "type": "video",
                                "format": "mp4",
                                "embedded_videos": []
                            })
                            
                            # 埋め込み動画を個別に抽出
                            if self.extract_embedded_videos_properly:
                                for i, video in enumerate(videos):
                                    embedded_video_path = self._extract_embedded_video_properly(
                                        pptx_path, slide_number, i + 1
                                    )
                                    if embedded_video_path:
                                        video['extracted_path'] = str(embedded_video_path)
                                        images_info[-1]["embedded_videos"].append(video)
                        else:
                            # 動画保存に失敗した場合はPNGとして保存
                            self.logger.warning(f"スライド {slide_number} の動画保存に失敗。PNGとして保存します。")
                            image_path = self._save_slide_as_image(pptx_path, slide_number)
                            if image_path:
                                images_info.append({
                                    "slide_number": slide_number,
                                    "image_path": str(image_path),
                                    "filename": image_path.name,
                                    "description": f"スライド {slide_number}",
                                    "type": "image",
                                    "format": "png",
                                    "embedded_videos": videos
                                })
                    else:
                        # 従来の処理方式（動画をGIFとして保存）
                        video_info = videos[0]  # 最初の動画を使用
                        
                        if video_info.get('is_embedded'):
                            # 埋め込み動画の場合、動画として扱う
                            video_path = video_info.get('video_path')
                            if video_path and Path(video_path).exists():
                                # 実際の動画ファイルがある場合
                                images_info.append({
                                    "slide_number": slide_number,
                                    "image_path": None,  # 動画のため画像パスはなし
                                    "filename": f"slide_{slide_number:02d}_video",
                                    "description": f"スライド {slide_number} の埋め込み動画",
                                    "type": "embedded_video",
                                    "embedded_videos": [video_info],
                                    "video_file": video_path
                                })
                                self.logger.info(f"スライド {slide_number} の埋め込み動画を検出")
                                self.logger.info(f"スライド {slide_number} の動画ファイルを保存: {video_path}")
                            else:
                                # 動画ファイルがない場合（メタデータのみ）
                                images_info.append({
                                    "slide_number": slide_number,
                                    "image_path": None,
                                    "filename": f"slide_{slide_number:02d}_video",
                                    "description": f"スライド {slide_number} の埋め込み動画（メタデータのみ）",
                                    "type": "embedded_video",
                                    "embedded_videos": [video_info]
                                })
                                self.logger.info(f"スライド {slide_number} の埋め込み動画を検出（メタデータのみ）")
                                self.logger.warning(f"スライド {slide_number} の動画ファイルが抽出できませんでした")
                        else:
                            # 外部動画ファイルの場合
                            gif_path = video_info.get('gif_path')
                            if gif_path and Path(gif_path).exists():
                                images_info.append({
                                    "slide_number": slide_number,
                                    "image_path": gif_path,
                                    "filename": Path(gif_path).name,
                                    "description": f"スライド {slide_number} の動画",
                                    "type": "video",
                                    "embedded_videos": [video_info]
                                })
                                self.logger.info(f"スライド {slide_number} の動画をGIFとして保存: {gif_path}")
                else:
                    # 動画がない場合は通常の画像を保存
                    image_path = self._save_slide_as_image(pptx_path, slide_number)
                    if image_path:
                        images_info.append({
                            "slide_number": slide_number,
                            "image_path": str(image_path),
                            "filename": image_path.name,
                            "description": f"スライド {slide_number}",
                            "type": "image",
                            "format": "png",
                            "embedded_videos": []
                        })
            
            self.logger.info(f"統合処理完了: {len(images_info)} スライド")
            return images_info
            
        except Exception as e:
            self.logger.error(f"統合処理でエラー: {e}")
            return self.extract_slides_as_images(pptx_path)
    
    def _save_slide_as_image(self, pptx_path: str, slide_number: int) -> Optional[Path]:
        """
        スライド全体を画像として保存
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            
        Returns:
            保存された画像のパス
        """
        try:
            # スライド全体を画像として抽出
            slide_image = self._extract_slide_as_image(pptx_path, slide_number)
            
            # 画像が見つからない場合はスキップ
            if slide_image is None:
                return None
            
            # 動画用にリサイズ
            resized_image = self._resize_for_video(slide_image)
            
            # ファイル名を生成
            filename = f"slide_{slide_number:02d}.{self.image_format}"
            output_path = self.output_dir / filename
            
            # 画像を保存
            if self.image_format.lower() == 'png':
                resized_image.save(output_path, 'PNG', optimize=True)
            else:
                resized_image.save(output_path, 'JPEG', quality=self.image_quality, optimize=True)
            
            self.logger.info(f"スライド {slide_number} の画像を保存しました: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"スライド {slide_number} の保存に失敗: {e}")
            return None
    
    def _extract_slide_as_image(self, pptx_path: str, slide_number: int) -> Optional[Image.Image]:
        """
        スライド全体を画像として抽出（python-pptx + PIL描画）
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            
        Returns:
            PIL画像オブジェクト
        """
        try:
            # スライドオブジェクトを取得
            prs = Presentation(pptx_path)
            if slide_number <= len(prs.slides):
                slide = prs.slides[slide_number - 1]
                
                # スライドを画像として描画
                slide_image = self._render_slide_to_image(slide)
                
                if slide_image:
                    self.logger.info(f"スライド {slide_number} 全体を画像化しました")
                    return slide_image
                else:
                    self.logger.warning(f"スライド {slide_number} の画像化に失敗しました")
                    return None
            else:
                self.logger.error(f"スライド {slide_number} が見つかりません")
                return None
            
        except Exception as e:
            self.logger.warning(f"スライド {slide_number} の画像抽出でエラー: {e}")
            return None
    
    def _render_slide_to_image(self, slide) -> Optional[Image.Image]:
        """
        スライドをPIL画像としてレンダリング
        
        Args:
            slide: PPTXスライドオブジェクト
            
        Returns:
            PIL画像オブジェクト
        """
        try:
            # スライドのサイズを取得（デフォルト値を使用）
            slide_width = 9144000  # 10インチ in EMU
            slide_height = 6858000  # 7.5インチ in EMU
            
            # スライドを画像としてレンダリング
            image = self._create_slide_image(slide, slide_width, slide_height)
            
            return image
            
        except Exception as e:
            self.logger.error(f"スライドのレンダリングに失敗: {e}")
            return None
    
    def _create_slide_image(self, slide, width: int, height: int) -> Image.Image:
        """
        スライドの内容を基に画像を作成（高解像度版）
        
        Args:
            slide: PPTXスライドオブジェクト
            width: スライド幅
            height: スライド高さ
            
        Returns:
            PIL画像オブジェクト
        """
        # 高解像度でスライドサイズをピクセルに変換（1920x1080相当）
        pixel_width = 1920
        pixel_height = 1080
        
        # 新しい画像を作成（高解像度）
        image = Image.new('RGB', (pixel_width, pixel_height), 'white')
        draw = ImageDraw.Draw(image)
        
        # 位置情報がない場合のフォールバック用
        y_offset = 100
        
        # 既に描画された領域を追跡（重複防止用）
        occupied_areas = []
        
        def is_area_occupied(x, y, width, height):
            """指定された領域が既に使用されているかチェック"""
            new_area = (x, y, x + width, y + height)
            for occupied in occupied_areas:
                # 重複チェック（矩形の重なり判定）
                if not (new_area[2] <= occupied[0] or new_area[0] >= occupied[2] or
                       new_area[3] <= occupied[1] or new_area[1] >= occupied[3]):
                    return True
            return False
        
        def add_occupied_area(x, y, width, height):
            """使用済み領域を追加"""
            occupied_areas.append((x, y, x + width, y + height))
        
        # フォントを設定（日本語対応）
        try:
            # 日本語フォントを優先的に試行
            font_paths = [
                "C:/Windows/Fonts/msgothic.ttc",      # MS Gothic
                "C:/Windows/Fonts/yu Gothic.ttc",     # Yu Gothic
                "C:/Windows/Fonts/meiryo.ttc",        # Meiryo
                "C:/Windows/Fonts/arial.ttf",         # Arial (フォールバック)
            ]
            
            title_font = None
            body_font = None
            small_font = None
            
            for font_path in font_paths:
                if os.path.exists(font_path):
                    try:
                        title_font = ImageFont.truetype(font_path, 72)  # タイトル用
                        body_font = ImageFont.truetype(font_path, 48)   # 本文用
                        small_font = ImageFont.truetype(font_path, 36)  # 小さいテキスト用
                        break
                    except:
                        continue
            
            # フォントが見つからない場合はデフォルトを使用
            if not title_font:
                title_font = ImageFont.load_default()
                body_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
                
        except Exception as e:
            self.logger.warning(f"フォント設定でエラー: {e}")
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
        
        # スライド内の要素を分類（改良版）
        text_shapes = []
        image_shapes = []
        
        for shape in slide.shapes:
            # テキスト要素の判定
            if hasattr(shape, "text") and shape.text.strip():
                text_shapes.append(shape)
            
            # 画像要素の判定（複数の方法で確認）
            is_image = False
            
            # Movieオブジェクト（shape_type == 18）を除外
            if hasattr(shape, 'shape_type') and shape.shape_type == 18:
                self.logger.debug(f"Movieオブジェクトを除外: {shape}")
                continue
            
            if hasattr(shape, 'image') and shape.image:
                is_image = True
            elif hasattr(shape, 'shape_type'):
                # 画像関連のシェイプタイプをチェック（Movieを除く）
                image_types = [13, 14, 15, 16, 17, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30]
                if shape.shape_type in image_types:
                    is_image = True
            elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                # 塗りつぶしタイプが画像の場合
                if shape.fill.type == 2:  # MSO_FILL.PICTURE
                    is_image = True
            
            if is_image:
                image_shapes.append(shape)
                self.logger.debug(f"画像要素を発見: {shape.shape_type if hasattr(shape, 'shape_type') else 'unknown'}")
        
        self.logger.debug(f"テキスト要素数: {len(text_shapes)}, 画像要素数: {len(image_shapes)}")
        
        # テキスト要素を上半分に配置
        text_area_height = pixel_height // 2
        text_y_offset = 50
        
        for shape in text_shapes:
            text = shape.text.strip()
            
            # "Agenda"を"今日のテーマ"に置き換え
            if "Agenda" in text:
                text = text.replace("Agenda", "今日のテーマ")
            
            # テキストの位置とサイズを取得（改良版）
            try:
                if hasattr(shape, 'left') and hasattr(shape, 'top') and hasattr(shape, 'width') and hasattr(shape, 'height'):
                    # EMUからピクセルに変換（より正確な計算）
                    slide_width_emu = slide.slide_width
                    slide_height_emu = slide.slide_height
                    
                    x = int(shape.left * pixel_width / slide_width_emu)
                    y = int(shape.top * pixel_height / slide_height_emu)
                    shape_width = int(shape.width * pixel_width / slide_width_emu)
                    shape_height = int(shape.height * pixel_height / slide_height_emu)
                    
                    # 上半分に制限
                    y = max(50, min(y, text_area_height - shape_height - 50))
                    x = max(50, min(x, pixel_width - shape_width - 50))
                else:
                    # 位置情報がない場合は上半分に配置
                    x = 100
                    y = text_y_offset
                    shape_width = pixel_width - 200
                    shape_height = 100
                    text_y_offset += 120
            except Exception as e:
                self.logger.debug(f"位置情報の取得でエラー: {e}")
                x = 100
                y = text_y_offset
                shape_width = pixel_width - 200
                shape_height = 100
                text_y_offset += 120
                
                # 重複チェック
                if is_area_occupied(x, y, shape_width, shape_height):
                    # 重複している場合は位置を調整
                    adjusted_x = x
                    adjusted_y = y
                    attempts = 0
                    while is_area_occupied(adjusted_x, adjusted_y, shape_width, shape_height) and attempts < 10:
                        adjusted_y += 50
                        attempts += 1
                    
                    if attempts >= 10:
                        # 調整できない場合はスキップ
                        self.logger.debug(f"スライド {slide.slide_id} のテキスト要素の位置調整に失敗")
                        continue
                    
                    x = adjusted_x
                    y = adjusted_y
                
                # テキストの長さに応じてフォントサイズを調整（改良版）
                if len(text) > 200:
                    font = small_font
                    max_width = shape_width - 20
                elif len(text) > 100:
                    font = body_font
                    max_width = shape_width - 20
                else:
                    font = title_font
                    max_width = shape_width - 20
                
                # テキストを描画（改行対応、改良版）
                lines = text.split('\n')
                current_y = y
                max_y = y + shape_height
                
                for line in lines:
                    if line.strip() and current_y < max_y - font.size:
                        # 長い行を分割
                        words = line.split()
                        current_line = ""
                        for word in words:
                            test_line = current_line + " " + word if current_line else word
                            if draw.textlength(test_line, font=font) < max_width:
                                current_line = test_line
                            else:
                                if current_line:
                                    draw.text((x + 10, current_y), current_line, fill='black', font=font)
                                    current_y += font.size + 5
                                current_line = word
                        
                        if current_line:
                            draw.text((x + 10, current_y), current_line, fill='black', font=font)
                            current_y += font.size + 10
                
                # 使用済み領域を追加
                add_occupied_area(x, y, shape_width, shape_height)
            
        # 画像要素を下半分に配置（さらに大きく表示）
        image_area_start = pixel_height // 2
        image_y_offset = image_area_start + 10  # 余白をさらに減らして画像を大きく
        
        for shape in image_shapes:
            try:
                # 画像データを取得（複数の方法で試行）
                image_data = None
                
                # 方法1: shape.image.blob
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_data = shape.image.blob
                        self.logger.debug(f"方法1で画像データを取得: {len(image_data)} bytes")
                    except Exception as e:
                        self.logger.debug(f"方法1で失敗: {e}")
                
                # 方法2: shape.fill.fore_color.rgb
                if not image_data and hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                    try:
                        image_data = shape.fill.fore_color.rgb
                        self.logger.debug(f"方法2で画像データを取得: {len(image_data)} bytes")
                    except Exception as e:
                        self.logger.debug(f"方法2で失敗: {e}")
                
                # 方法3: その他の属性から取得
                if not image_data:
                    for attr_name in ['blob', 'image_data', 'data']:
                        if hasattr(shape, attr_name):
                            try:
                                image_data = getattr(shape, attr_name)
                                self.logger.debug(f"方法3で画像データを取得: {len(image_data)} bytes")
                                break
                            except Exception as e:
                                self.logger.debug(f"方法3で失敗: {e}")
                
                if not image_data:
                    # Movieオブジェクトの場合はスキップ
                    if hasattr(shape, 'shape_type') and shape.shape_type == 18:  # 動画シェイプ
                        self.logger.debug(f"動画シェイプをスキップ: {shape}")
                        continue
                    else:
                        self.logger.warning(f"画像データを取得できませんでした: {shape}")
                        continue
                
                pil_image = Image.open(io.BytesIO(image_data))
                self.logger.debug(f"画像を読み込みました: {pil_image.size}")
                
                # 画像の位置とサイズを取得（改良版）
                try:
                    if hasattr(shape, 'left') and hasattr(shape, 'top') and hasattr(shape, 'width') and hasattr(shape, 'height'):
                        # EMUからピクセルに変換（より正確な計算）
                        slide_width_emu = slide.slide_width
                        slide_height_emu = slide.slide_height
                        
                        x = int(shape.left * pixel_width / slide_width_emu)
                        y = int(shape.top * pixel_height / slide_height_emu)
                        shape_width = int(shape.width * pixel_width / slide_width_emu)
                        shape_height = int(shape.height * pixel_height / slide_height_emu)
                        
                        # 下半分に制限（さらに大きく表示）
                        x = max(20, min(x, pixel_width - 40))  # 左右の余白をさらに減らす
                        y = max(image_area_start + 10, min(y, pixel_height - 20))  # 上下の余白をさらに減らす
                        shape_width = min(shape_width, pixel_width - x - 20)
                        shape_height = min(shape_height, pixel_height - y - 20)
                    else:
                        # 位置情報がない場合は下半分に配置（さらに大きく表示）
                        x = (pixel_width - pil_image.width) // 2
                        y = image_y_offset
                        shape_width = min(pil_image.width, pixel_width - 40)  # 左右の余白をさらに減らす
                        shape_height = min(pil_image.height, pixel_height - image_y_offset - 20)  # 上下の余白をさらに減らす
                        image_y_offset += shape_height + 20  # 画像間の余白もさらに減らす
                except Exception as e:
                    self.logger.debug(f"画像位置情報の取得でエラー: {e}")
                    x = (pixel_width - pil_image.width) // 2
                    y = image_y_offset
                    shape_width = min(pil_image.width, pixel_width - 40)  # 左右の余白をさらに減らす
                    shape_height = min(pil_image.height, pixel_height - image_y_offset - 20)  # 上下の余白をさらに減らす
                    image_y_offset += shape_height + 20  # 画像間の余白もさらに減らす
                
                # 重複チェック
                if is_area_occupied(x, y, shape_width, shape_height):
                    # 重複している場合は位置を調整
                    adjusted_x = x
                    adjusted_y = y
                    attempts = 0
                    while is_area_occupied(adjusted_x, adjusted_y, shape_width, shape_height) and attempts < 10:
                        adjusted_y += 50
                        attempts += 1
                    
                    if attempts >= 10:
                        # 調整できない場合はスキップ
                        self.logger.debug(f"スライド {slide.slide_id} の画像要素の位置調整に失敗")
                        continue
                    
                    x = adjusted_x
                    y = adjusted_y
                
                # 下半分の領域を最大限活用して画像をリサイズ
                # 下半分の利用可能な領域を計算（より大きな余白を確保）
                available_width = pixel_width - 60  # 左右の余白を増加
                available_height = pixel_height - image_area_start - 40  # 下半分の高さ（上下の余白を増加）
                
                # 画像の縦横比を計算
                image_ratio = pil_image.width / pil_image.height
                available_ratio = available_width / available_height
                
                # 縦横比に応じて最適なサイズを計算
                if image_ratio > available_ratio:
                    # 横長画像：幅を最大限活用
                    new_width = available_width
                    new_height = int(available_width / image_ratio)
                else:
                    # 縦長画像：高さを最大限活用
                    new_height = available_height
                    new_width = int(available_height * image_ratio)
                
                # 画像をリサイズ
                pil_image = pil_image.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                # 下半分の中央に配置
                # x位置は中央に配置
                original_x = (pixel_width - new_width) // 2
                # y位置は下半分の中央に配置
                original_y = image_area_start + (pixel_height - image_area_start - new_height) // 2
                image.paste(pil_image, (original_x, original_y))
                
                # 使用済み領域を追加
                add_occupied_area(original_x, original_y, new_width, new_height)
                
            except Exception as e:
                self.logger.debug(f"画像の処理に失敗: {e}")
                continue
        
        # 画像が見つからない場合のフォールバック処理
        if len(image_shapes) == 0:
            self.logger.debug(f"スライド {slide.slide_id} に画像要素が見つかりませんでした")
            # すべてのシェイプをログ出力してデバッグ
            for i, shape in enumerate(slide.shapes):
                shape_info = f"Shape {i}: type={getattr(shape, 'shape_type', 'unknown')}"
                if hasattr(shape, 'text') and shape.text.strip():
                    shape_info += f", text='{shape.text[:50]}...'"
                if hasattr(shape, 'image'):
                    shape_info += f", has_image={shape.image is not None}"
                self.logger.debug(shape_info)
        
        return image
    
    def _resize_for_video(self, image: Image.Image) -> Image.Image:
        """
        画像を動画用にリサイズ
        
        Args:
            image: 元の画像
            
        Returns:
            リサイズされた画像
        """
        # アスペクト比を保持してリサイズ
        image.thumbnail((self.width, self.height), Image.Resampling.LANCZOS)
        
        # 背景を追加して指定サイズに
        new_image = Image.new('RGB', (self.width, self.height), 'white')
        
        # 画像を中央に配置
        x = (self.width - image.width) // 2
        y = (self.height - image.height) // 2
        new_image.paste(image, (x, y))
        
        return new_image
    
    def process_images_for_script(self, pptx_path: str, script_data: Dict[str, Any], videos_info: Optional[List[Dict[str, Any]]] = None) -> Dict[str, Any]:
        """
        台本に基づいて画像を処理
        
        Args:
            pptx_path: PPTXファイルのパス
            script_data: 台本データ
            videos_info: 既に抽出済みの動画情報（オプション）
            
        Returns:
            処理された画像情報
        """
        self.logger.info("台本に基づいて画像を処理中...")
        
        # 既に動画確認済みの場合は、その情報を使用
        if videos_info:
            self.logger.info("既に動画確認済みの情報を使用します")
            images_info = videos_info
        else:
            # 動画確認済みでない場合は、通常の画像抽出
            self.logger.info("通常の画像抽出を実行します")
            images_info = self.extract_slides_as_images(pptx_path)
        
        # 台本の対話に基づいて画像のタイミングを設定
        processed_images = self._assign_timing_to_images(images_info, script_data)
        
        # メタデータを保存
        metadata_path = self.output_dir / "images_metadata.json"
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(processed_images, f, ensure_ascii=False, indent=2)
        
        self.logger.info(f"画像処理完了: {len(processed_images)} 画像")
        return processed_images
    
    def _assign_timing_to_images(self, images_info: List[Dict[str, Any]], script_data: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        画像にタイミング情報を割り当て
        
        Args:
            images_info: 画像情報のリスト
            script_data: 台本データ
            
        Returns:
            タイミング情報付きの画像リスト
        """
        processed_images = []
        
        # 台本の対話からスライドファイル名とタイミングを抽出
        slide_timings = {}
        
        for dialogue in script_data.get("dialogue", []):
            slide_file = dialogue.get("slide_file", "slide_01.png")
            timestamp = dialogue.get("timestamp", "00:00:00")
            
            if slide_file not in slide_timings:
                slide_timings[slide_file] = {
                    "start_time": timestamp,
                    "end_time": timestamp
                }
            else:
                slide_timings[slide_file]["end_time"] = timestamp
        
        # 画像情報にタイミングを追加
        for image_info in images_info:
            slide_file = image_info["filename"]  # 実際のファイル名を使用
            timing = slide_timings.get(slide_file, {
                "start_time": "00:00:00",
                "end_time": "00:00:00"
            })
            
            processed_image = {
                **image_info,
                "start_time": timing["start_time"],
                "end_time": timing["end_time"],
                "duration": self._calculate_duration(timing["start_time"], timing["end_time"])
            }
            
            processed_images.append(processed_image)
        
        return processed_images
    
    def _calculate_duration(self, start_time: str, end_time: str) -> int:
        """
        時間の差を秒数で計算
        
        Args:
            start_time: 開始時間（HH:MM:SS形式）
            end_time: 終了時間（HH:MM:SS形式）
            
        Returns:
            秒数
        """
        def time_to_seconds(time_str: str) -> int:
            parts = time_str.split(':')
            return int(parts[0]) * 3600 + int(parts[1]) * 60 + int(parts[2])
        
        start_seconds = time_to_seconds(start_time)
        end_seconds = time_to_seconds(end_time)
        
        return max(0, end_seconds - start_seconds)
    
    def create_transition_effects(self, images_info: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        トランジション効果を追加
        
        Args:
            images_info: 画像情報のリスト
            
        Returns:
            トランジション効果付きの画像リスト
        """
        self.logger.info("トランジション効果を追加中...")
        
        processed_images = []
        
        for i, image_info in enumerate(images_info):
            processed_image = {
                **image_info,
                "transition": {
                    "type": "fade",
                    "duration": 1.0,  # 1秒のフェード
                    "direction": "in" if i == 0 else "cross"
                }
            }
            
            processed_images.append(processed_image)
        
        self.logger.info("トランジション効果の追加完了")
        return processed_images
    
    def save_images_metadata(self, images_info: List[Dict[str, Any]], output_path: str):
        """
        画像メタデータを保存
        
        Args:
            images_info: 画像情報のリスト
            output_path: 出力ファイルパス
        """
        try:
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(images_info, f, ensure_ascii=False, indent=2)
            
            self.logger.info(f"画像メタデータを保存しました: {output_path}")
            
        except Exception as e:
            self.logger.error(f"画像メタデータの保存に失敗: {e}")
            raise 

    def _can_use_win32com(self) -> bool:
        """win32comが使用可能かチェック"""
        import platform
        if platform.system() != 'Windows':
            return False
        
        try:
            import win32com.client
            return True
        except ImportError:
            return False

    def _can_use_libreoffice(self) -> bool:
        """LibreOfficeが使用可能かチェック"""
        import shutil
        
        libreoffice_paths = [
            "soffice",
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        ]
        
        for path in libreoffice_paths:
            if shutil.which(path) or Path(path).exists():
                # pdf2imageも必要
                try:
                    import pdf2image
                    return True
                except ImportError:
                    return False
        
        return False

    def _extract_with_powerpoint(self, pptx_path: str) -> Optional[List[Dict[str, Any]]]:
        """PowerPoint (win32com) を使用して画像を抽出"""
        try:
            import win32com.client
            import pythoncom
            import os
            
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            self.logger.info(f"PowerPointアプリケーション初期化完了: {powerpoint.Version}")
            try:
                powerpoint.Visible = 0
                self.logger.info("PowerPoint Visible設定完了")
            except Exception as visible_error:
                self.logger.warning(f"Visible設定でエラーが発生しましたが無視して続行: {visible_error}")
                pass
            
            # 出力ディレクトリの存在を確認
            if not self.output_dir.exists():
                self.output_dir.mkdir(parents=True, exist_ok=True)
                self.logger.info(f"出力ディレクトリを作成: {self.output_dir}")
            
            try:
                pptx_abs_path = str(Path(pptx_path).resolve())
                self.logger.info(f"PowerPointプレゼンテーションを開いています: {pptx_abs_path}")
                presentation = powerpoint.Presentations.Open(pptx_abs_path, WithWindow=False)
                self.logger.info(f"プレゼンテーション読み込み完了: {presentation.Slides.Count} スライド")
                
                images_info = []
                total_slides = presentation.Slides.Count
                
                # 設定からエクスポート品質を取得
                export_width = config.get("image.export_width", 1920)
                export_height = config.get("image.export_height", 1080)
                
                for i in range(1, total_slides + 1):
                    self.logger.info(f"=== スライド {i}/{total_slides} の処理開始 ===")
                    
                    filename = f"slide_{i:02d}.png"
                    output_path = self.output_dir / filename
                    
                    # 高品質でエクスポート
                    try:
                        # 出力ディレクトリが存在することを確認
                        output_path.parent.mkdir(parents=True, exist_ok=True)
                        
                        # 絶対パスを使用
                        abs_output_path = str(output_path.resolve())
                        self.logger.info(f"スライド {i} をエクスポート中: {abs_output_path}")
                        
                        # スライドオブジェクトを取得
                        slide = presentation.Slides(i)
                        self.logger.info(f"スライド {i} オブジェクト取得完了")
                        
                        # エクスポート実行
                        slide.Export(abs_output_path, "PNG", export_width, export_height)
                        self.logger.info(f"スライド {i} のExport()呼び出し完了")

                        # ファイルが実際に作成されたか確認（最大5秒待つ）
                        import time
                        wait_start = time.time()
                        while not output_path.exists():
                            if time.time() - wait_start > 5.0:
                                raise Exception(f"エクスポートされたファイルが5秒以内に見つかりません: {abs_output_path}")
                            time.sleep(0.1)
                        
                        self.logger.info(f"スライド {i} のエクスポート成功: {abs_output_path}")
                        
                        # 画像情報をリストに追加
                        images_info.append({
                            "slide_number": i,
                            "image_path": str(output_path),
                            "filename": filename,
                            "description": f"スライド {i}"
                        })
                        self.logger.info(f"スライド {i} の画像情報をリストに追加完了")
                            
                    except Exception as export_error:
                        self.logger.warning(f"スライド {i} のエクスポートでエラー: {export_error}")
                        self.logger.warning(f"エラーの詳細: {type(export_error).__name__}: {str(export_error)}")
                        # デフォルトサイズで再試行
                        try:
                            self.logger.info(f"スライド {i} をデフォルトサイズで再試行中...")
                            slide.Export(abs_output_path, "PNG")
                            self.logger.info(f"スライド {i} のデフォルトExport()呼び出し完了")
                            
                            if not output_path.exists():
                                raise Exception(f"デフォルトエクスポートでもファイルが見つかりません: {abs_output_path}")
                            self.logger.info(f"スライド {i} のデフォルトエクスポート成功: {abs_output_path}")
                            
                            # 画像情報をリストに追加
                            images_info.append({
                                "slide_number": i,
                                "image_path": str(output_path),
                                "filename": filename,
                                "description": f"スライド {i}"
                            })
                            self.logger.info(f"スライド {i} の画像情報をリストに追加完了（デフォルト）")
                            
                        except Exception as retry_error:
                            self.logger.error(f"スライド {i} のエクスポートに完全に失敗: {retry_error}")
                            self.logger.error(f"再試行エラーの詳細: {type(retry_error).__name__}: {str(retry_error)}")
                            # このスライドはスキップして次に進む
                            pass
                    
                                    # ループの外側で画像情報を追加する処理は削除（ループ内で追加済み）
                
                presentation.Close()
                self.logger.info(f"=== PowerPoint COM Export 処理完了 ===")
                self.logger.info(f"処理対象スライド数: {total_slides}")
                self.logger.info(f"成功したスライド数: {len(images_info)}")
                self.logger.info(f"生成された画像ファイル: {[info['filename'] for info in images_info]}")
                
                if len(images_info) == 0:
                    self.logger.warning("PowerPoint COM Export: 画像が生成されませんでした")
                    return None
                elif len(images_info) < total_slides:
                    self.logger.warning(f"PowerPoint COM Export: 一部のスライドでエラーが発生しました ({len(images_info)}/{total_slides})")
                
                return images_info
                
            finally:
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
        except Exception as e:
            self.logger.error(f"PowerPoint経由での画像抽出に失敗: {e}")
            return None

    def _extract_with_powerpoint_saveas(self, pptx_path: str) -> Optional[List[Dict[str, Any]]]:
        """PowerPoint SaveAs機能を使用して画像を抽出"""
        try:
            import win32com.client
            import pythoncom
            import tempfile
            import shutil
            
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            
            # 一時ディレクトリを作成
            temp_dir = Path(tempfile.mkdtemp())
            
            try:
                # プレゼンテーションを開く
                pptx_abs_path = str(Path(pptx_path).resolve())
                self.logger.info(f"PowerPoint SaveAs: プレゼンテーションを開いています: {pptx_abs_path}")
                presentation = powerpoint.Presentations.Open(pptx_abs_path, WithWindow=False)
                self.logger.info(f"PowerPoint SaveAs: プレゼンテーション読み込み完了: {presentation.Slides.Count} スライド")
                
                # PNG形式で保存（全スライドが個別のPNGファイルとして保存される）
                output_folder = temp_dir / "slides"
                output_folder.mkdir(exist_ok=True)
                
                # SaveAsで画像として保存
                # ppSaveAsPNG = 18
                save_path = str(output_folder / "slide")
                self.logger.info(f"PowerPoint SaveAs: PNG形式で保存中: {save_path}")
                presentation.SaveAs(save_path, 18)
                
                # 生成された画像ファイルを収集
                images_info = []
                png_files = sorted(output_folder.glob("slide*.PNG"))
                
                self.logger.info(f"PowerPoint SaveAs: {len(png_files)} 個のPNGファイルを発見")
                
                for i, png_file in enumerate(png_files, 1):
                    # 出力先にコピー
                    filename = f"slide_{i:02d}.png"
                    output_path = self.output_dir / filename
                    
                    # 画像をリサイズして保存
                    from PIL import Image
                    img = Image.open(png_file)
                    img = img.resize((self.width, self.height), Image.Resampling.LANCZOS)
                    img.save(output_path, 'PNG', optimize=True)
                    
                    images_info.append({
                        "slide_number": i,
                        "image_path": str(output_path),
                        "filename": filename,
                        "description": f"スライド {i}"
                    })
                    
                    self.logger.info(f"PowerPoint SaveAs: スライド {i} を画像化完了")
                
                presentation.Close()
                
                if images_info:
                    self.logger.info(f"PowerPoint SaveAsで {len(images_info)} 枚の画像を生成")
                    return images_info
                else:
                    self.logger.warning("PowerPoint SaveAs: 画像ファイルが生成されませんでした")
                    return None
                    
            finally:
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                # 一時ディレクトリを削除
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
                
        except Exception as e:
            self.logger.error(f"PowerPoint SaveAs経由での画像抽出に失敗: {e}")
            return None

    def _extract_with_libreoffice(self, pptx_path: str) -> Optional[List[Dict[str, Any]]]:
        """LibreOfficeを使用して画像を抽出"""
        try:
            import subprocess
            import tempfile
            import pdf2image
            
            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as temp_dir:
                # LibreOfficeでPDFに変換
                cmd = [
                    'soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, pptx_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    self.logger.error(f"LibreOffice PDF変換に失敗: {result.stderr}")
                    return None
                
                # PDFファイルのパスを取得
                pdf_path = Path(temp_dir) / f"{Path(pptx_path).stem}.pdf"
                if not pdf_path.exists():
                    self.logger.error("PDFファイルが生成されませんでした")
                    return None
                
                # PDFから画像に変換
                images = pdf2image.convert_from_path(
                    str(pdf_path),
                    dpi=300,  # 高解像度
                    size=(1920, 1080)  # 動画用サイズ
                )
                
                images_info = []
                for i, image in enumerate(images, 1):
                    filename = f"slide_{i:02d}.png"
                    output_path = self.output_dir / filename
                    
                    # 画像を保存
                    image.save(output_path, 'PNG', optimize=True)
                    
                    images_info.append({
                        "slide_number": i,
                        "image_path": str(output_path),
                        "filename": filename,
                        "description": f"スライド {i}"
                    })
                
                self.logger.info(f"LibreOfficeで {len(images_info)} 枚のスライド画像を生成完了")
                return images_info
                
        except Exception as e:
            self.logger.error(f"LibreOffice経由での画像抽出に失敗: {e}")
            return None

    def _extract_slides_with_pil(self, pptx_path: str) -> List[Dict[str, Any]]:
        """標準的な方法（PIL描画）で画像を抽出（フォールバック）"""
        try:
            prs = Presentation(pptx_path)
            images_info = []
            total_slides = len(prs.slides)
            
            for i, slide in enumerate(prs.slides, 1):
                self.logger.info(f"スライド {i}/{total_slides} を画像化中...")
                
                # スライド全体を画像として保存
                image_path = self._save_slide_as_image(pptx_path, i)
                
                if image_path:
                    images_info.append({
                        "slide_number": i,
                        "image_path": str(image_path),
                        "filename": image_path.name,
                        "description": f"スライド {i}"
                    })
                    self.logger.debug(f"スライド {i} を画像化: {image_path}")
            
            self.logger.info(f"合計 {len(images_info)} 枚のスライド画像を生成完了（全{total_slides}スライド）")
            return images_info
            
        except Exception as e:
            self.logger.error(f"画像抽出に失敗: {e}")
            raise 

    def extract_embedded_videos_from_slide(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """
        スライドから埋め込み動画を抽出
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            
        Returns:
            動画情報のリスト
        """
        if not self.extract_embedded_videos:
            return []
        
        try:
            # PowerPoint COMを優先して使用
            if self._can_use_win32com():
                self.logger.debug(f"PowerPoint COMでスライド {slide_number} の動画を検索中...")
                videos = self._extract_videos_with_powerpoint(pptx_path, slide_number)
                if videos:
                    self.logger.info(f"PowerPoint COMで動画を検出: {len(videos)}個")
                    return videos
                else:
                    self.logger.debug(f"PowerPoint COMで動画を検出できませんでした")
            
            # フォールバック: python-pptxを使用
            self.logger.debug(f"python-pptxでスライド {slide_number} の動画を検索中...")
            from pptx import Presentation
            import tempfile
            import subprocess
            import shutil
            
            prs = Presentation(pptx_path)
            if slide_number > len(prs.slides):
                self.logger.warning(f"スライド {slide_number} は存在しません")
                return []
            
            slide = prs.slides[slide_number - 1]
            videos_info = []
            
            # スライド内の動画を検索（強化版）
            for shape in slide.shapes:
                self.logger.debug(f"シェイプチェック: {shape.shape_type}, 名前: {shape.name}")
                
                # 動画の可能性があるシェイプタイプ
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                video_types = [
                    MSO_SHAPE_TYPE.MEDIA,  # メディアオブジェクト
                ]
                
                # 1. shape_typeでチェック
                if hasattr(shape, 'shape_type') and shape.shape_type in video_types:
                    self.logger.info(f"動画候補シェイプ検出: {shape.shape_type}")
                    video_info = self._create_video_info_from_shape(shape, slide_number)
                    if video_info:
                        videos_info.append(video_info)
                
                # 2. メディア関連の属性をチェック
                elif hasattr(shape, 'media') or hasattr(shape, 'movie'):
                    self.logger.info(f"メディア属性を持つシェイプ検出")
                    video_info = self._create_video_info_from_shape(shape, slide_number)
                    if video_info:
                        videos_info.append(video_info)
                
                # 3. element内のvideoタグをチェック（より厳密に）
                elif hasattr(shape, 'element'):
                    # テキストボックスや画像は除外
                    if hasattr(shape, 'shape_type'):
                        if shape.shape_type == 17:  # TEXT_BOX
                            continue
                        elif shape.shape_type == 13:  # PICTURE
                            continue
                    
                    if self._check_element_for_video(shape.element):
                        self.logger.info(f"element内に動画タグ検出")
                        video_info = self._create_video_info_from_shape(shape, slide_number)
                        if video_info:
                            videos_info.append(video_info)
                
                # 4. 従来の方法（フォールバック）
                elif hasattr(shape, 'movie') and shape.movie:
                    try:
                        video_info = self._extract_video_from_shape(shape, slide_number)
                        if video_info:
                            videos_info.append(video_info)
                    except Exception as e:
                        self.logger.error(f"動画抽出エラー: {e}")
            
            if videos_info:
                self.logger.info(f"python-pptxで動画を検出: {len(videos_info)}個")
                return videos_info
            
            # 5. PPTXファイル構造の直接検査
            self.logger.debug(f"PPTXファイル構造でスライド {slide_number} の動画を検索中...")
            videos_info = self._extract_videos_from_pptx_structure(pptx_path, slide_number)
            if videos_info:
                self.logger.info(f"PPTXファイル構造で動画を検出: {len(videos_info)}個")
                return videos_info
            
            self.logger.info(f"スライド {slide_number} から {len(videos_info)} 個の動画を抽出")
            return videos_info
            
        except Exception as e:
            self.logger.error(f"スライド {slide_number} の動画抽出に失敗: {e}")
            return []
    
    def _extract_video_from_shape(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """
        シェイプから動画を抽出
        
        Args:
            shape: 動画を含むシェイプ
            slide_number: スライド番号
            
        Returns:
            動画情報の辞書
        """
        try:
            # 動画ファイルのパスを取得
            video_path = shape.movie.filename
            # 動画情報を取得
            if video_path and Path(video_path).exists():
                # 外部動画ファイルの場合
                video_info = self._get_video_info(video_path)
                if not video_info:
                    return None
                
                # 動画をGIFに変換
                gif_path = self._convert_video_to_gif(video_path, slide_number)
                if not gif_path:
                    return None
                
                return {
                    "slide_number": slide_number,
                    "original_video_path": video_path,
                    "gif_path": str(gif_path),
                    "duration": video_info.get("duration", 0),
                    "width": video_info.get("width", 0),
                    "height": video_info.get("height", 0),
                    "filename": gif_path.name,
                    "description": f"スライド {slide_number} の埋め込み動画"
                }
            else:
                # 埋め込み動画の場合、動画情報のみを返す
                self.logger.info(f"スライド {slide_number} に埋め込み動画を検出")
                return {
                    "slide_number": slide_number,
                    "original_video_path": None,
                    "gif_path": None,
                    "duration": 10,  # デフォルト10秒
                    "width": 640,
                    "height": 480,
                    "filename": f"slide_{slide_number:02d}_embedded_video",
                    "description": f"スライド {slide_number} の埋め込み動画",
                    "is_embedded": True
                }
            
        except Exception as e:
            self.logger.error(f"動画抽出エラー: {e}")
            return None
    
    def _get_video_info(self, video_path: str) -> Optional[Dict[str, Any]]:
        """
        動画ファイルの情報を取得
        
        Args:
            video_path: 動画ファイルのパス
            
        Returns:
            動画情報の辞書
        """
        try:
            import subprocess
            import json
            
            # FFprobeで動画情報を取得
            cmd = [
                'ffprobe', '-v', 'quiet', '-print_format', 'json',
                '-show_format', '-show_streams', video_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                self.logger.error(f"FFprobeエラー: {result.stderr}")
                return None
            
            info = json.loads(result.stdout)
            
            # 動画ストリームを取得
            video_stream = None
            for stream in info.get('streams', []):
                if stream.get('codec_type') == 'video':
                    video_stream = stream
                    break
            
            if not video_stream:
                self.logger.error("動画ストリームが見つかりません")
                return None
            
            # 動画情報を抽出
            duration = float(info.get('format', {}).get('duration', 0))
            width = int(video_stream.get('width', 0))
            height = int(video_stream.get('height', 0))
            
            # 最大動画長をチェック
            if duration > self.max_video_duration:
                self.logger.warning(f"動画が長すぎます: {duration}秒 > {self.max_video_duration}秒")
                duration = self.max_video_duration
            
            return {
                "duration": duration,
                "width": width,
                "height": height,
                "format": info.get('format', {}).get('format_name', 'unknown')
            }
            
        except Exception as e:
            self.logger.error(f"動画情報取得エラー: {e}")
            return None
    
    def _convert_video_to_gif(self, video_path: str, slide_number: int) -> Optional[Path]:
        """
        動画をGIFに変換
        
        Args:
            video_path: 動画ファイルのパス
            slide_number: スライド番号
            
        Returns:
            GIFファイルのパス
        """
        try:
            import subprocess
            
            # 出力ファイル名を生成
            filename = f"slide_{slide_number:02d}_video.gif"
            output_path = self.output_dir / filename
            
            # FFmpegでGIFに変換
            cmd = [
                'ffmpeg', '-y',
                '-i', video_path,
                '-vf', f'fps={self.gif_fps},scale={self.video_width}:{self.video_height}',
                '-t', str(self.max_video_duration),
                str(output_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                self.logger.error(f"GIF変換エラー: {result.stderr}")
                return None
            
            self.logger.info(f"動画をGIFに変換完了: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"GIF変換エラー: {e}")
            return None
    
    def process_slides_with_videos(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        スライドを処理し、動画も含めて画像情報を取得
        
        Args:
            pptx_path: PPTXファイルのパス
            
        Returns:
            画像・動画情報のリスト
        """
        # 通常の画像抽出
        images_info = self.extract_slides_as_images(pptx_path)
        
        # 動画抽出を追加
        for image_info in images_info:
            slide_number = image_info["slide_number"]
            videos = self.extract_embedded_videos_from_slide(pptx_path, slide_number)
            
            if videos:
                image_info["embedded_videos"] = videos
                self.logger.info(f"スライド {slide_number} に {len(videos)} 個の動画を追加")
        
        return images_info 

    def _extract_videos_with_powerpoint(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """
        PowerPoint COMを使用して動画を抽出
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            
        Returns:
            動画情報のリスト
        """
        try:
            import win32com.client
            import pythoncom
            import tempfile
            import shutil
            
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            
            try:
                # プレゼンテーションを開く
                pptx_abs_path = str(Path(pptx_path).resolve())
                presentation = powerpoint.Presentations.Open(pptx_abs_path, WithWindow=False)
                
                # スライドオブジェクトを取得
                slide = presentation.Slides(slide_number)
                
                # スライド内の動画を検索
                videos_info = []
                self.logger.info(f"スライド {slide_number} のシェイプ数を確認: {slide.Shapes.Count}")
                
                for i in range(1, slide.Shapes.Count + 1):
                    try:
                        shape = slide.Shapes(i)
                        self.logger.debug(f"シェイプ {i}: タイプ={shape.Type}, 名前={shape.Name}")
                        
                        # 動画オブジェクトかどうかを確認（複数の方法）
                        is_video = False
                        
                        # 方法1: MediaType
                        if hasattr(shape, 'MediaType'):
                            media_type = shape.MediaType
                            self.logger.debug(f"シェイプ {i} のメディアタイプ: {media_type}")
                            if media_type == 2:  # ppMediaTypeMovie
                                is_video = True
                        
                        # 方法2: シェイプ名で確認
                        if not is_video and hasattr(shape, 'Name'):
                            shape_name = shape.Name.lower()
                            if any(keyword in shape_name for keyword in ['movie', 'video', 'media']):
                                self.logger.debug(f"シェイプ {i} が動画として認識されました（名前ベース）")
                                is_video = True
                        
                        # 方法3: シェイプタイプで確認
                        if not is_video and hasattr(shape, 'Type'):
                            shape_type = shape.Type
                            # 動画関連のシェイプタイプ（PowerPointの定数）
                            # msoMedia = 16, その他の動画関連タイプ
                            if shape_type in [16, 18, 19, 20]:  # 動画関連のタイプ
                                self.logger.debug(f"シェイプ {i} が動画として認識されました（タイプベース）")
                                is_video = True
                        
                        # 方法4: PlaceholderFormatのチェック
                        if not is_video and hasattr(shape, 'PlaceholderFormat'):
                            try:
                                placeholder = shape.PlaceholderFormat
                                # ppPlaceholderMedia = 15
                                if hasattr(placeholder, 'Type') and placeholder.Type == 15:
                                    self.logger.debug(f"シェイプ {i} が動画として認識されました（PlaceholderFormat）")
                                    is_video = True
                            except Exception as e:
                                self.logger.debug(f"PlaceholderFormat確認でエラー: {e}")
                        
                        if is_video:
                            self.logger.info(f"スライド {slide_number} で動画を発見: シェイプ {i}")
                            # 実際の動画ファイルをMP4として抽出
                            video_info = self._extract_embedded_video_to_mp4(shape, slide_number)
                            if video_info:
                                videos_info.append(video_info)
                                self.logger.info(f"動画ファイルを抽出しました: {video_info['video_path']}")
                            else:
                                # フォールバック: メタデータのみ
                                video_info = self._extract_video_from_powerpoint_shape(shape, slide_number)
                                if video_info:
                                    videos_info.append(video_info)
                                self.logger.info(f"動画抽出成功: {video_info.get('filename', 'unknown')}")
                        else:
                            self.logger.debug(f"シェイプ {i} は動画ではありません")
                            
                    except Exception as e:
                        self.logger.debug(f"シェイプ {i} の動画確認でエラー: {e}")
                        continue
                
                presentation.Close()
                return videos_info
                
            finally:
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
        except Exception as e:
            self.logger.error(f"PowerPoint COM動画抽出でエラー: {e}")
            return []
    
    def _extract_video_from_powerpoint_shape(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """
        PowerPoint COMシェイプから動画を抽出
        
        Args:
            shape: PowerPoint COMシェイプオブジェクト
            slide_number: スライド番号
            
        Returns:
            動画情報の辞書
        """
        try:
            # 動画ファイルのパスを取得
            self.logger.debug(f"動画シェイプの詳細確認中...")
            
            # 複数の方法で動画パスを取得
            video_path = None
            
            try:
                # 方法1: MediaFormat.SourceFullName
                if hasattr(shape, 'MediaFormat'):
                    video_path = shape.MediaFormat.SourceFullName
                    self.logger.debug(f"MediaFormat.SourceFullName: {video_path}")
            except Exception as e:
                self.logger.debug(f"MediaFormat.SourceFullName取得エラー: {e}")
            
            # 方法2: 代替手段
            if not video_path:
                try:
                    # 動画ファイル名を取得
                    if hasattr(shape, 'Name'):
                        shape_name = shape.Name
                        self.logger.debug(f"シェイプ名: {shape_name}")
                        
                        # 埋め込み動画の場合、動画情報を取得
                        if "movie" in shape_name.lower() or "video" in shape_name.lower():
                            # 動画情報を取得
                            video_info = self._extract_embedded_video_to_temp(shape, slide_number)
                            if video_info:
                                self.logger.debug(f"埋め込み動画情報を取得: {video_info}")
                                return video_info
                except Exception as e:
                    self.logger.debug(f"代替動画パス取得エラー: {e}")
            
            if not video_path or not Path(video_path).exists():
                self.logger.warning(f"動画ファイルが見つかりません: {video_path}")
                return None
            
            # 動画情報を取得
            video_info = self._get_video_info(video_path)
            if not video_info:
                return None
            
            # 動画をGIFに変換
            gif_path = self._convert_video_to_gif(video_path, slide_number)
            if not gif_path:
                return None
            
            return {
                "slide_number": slide_number,
                "original_video_path": video_path,
                "gif_path": str(gif_path),
                "duration": video_info.get("duration", 0),
                "width": video_info.get("width", 0),
                "height": video_info.get("height", 0),
                "filename": gif_path.name,
                "description": f"スライド {slide_number} の埋め込み動画"
            }
            
        except Exception as e:
            self.logger.error(f"PowerPoint COM動画抽出エラー: {e}")
            return None 
    
    def _extract_embedded_video_to_temp(self, shape, slide_number: int) -> Optional[str]:
        """
        埋め込み動画の情報を取得（実際のファイルは作成しない）
        
        Args:
            shape: PowerPoint COMシェイプオブジェクト
            slide_number: スライド番号
            
        Returns:
            動画情報の辞書
        """
        try:
            # 動画のプロパティを取得
            if hasattr(shape, 'MediaFormat'):
                media_format = shape.MediaFormat
                
                # 動画の長さを取得
                duration = getattr(media_format, 'Duration', 10)  # デフォルト10秒
                self.logger.debug(f"動画の長さ: {duration}秒")
                
                # 動画のサイズを取得
                width = getattr(media_format, 'Width', 640)
                height = getattr(media_format, 'Height', 480)
                self.logger.debug(f"動画のサイズ: {width}x{height}")
                
                # 動画情報を返す（実際のファイルは作成しない）
                return {
                    "duration": duration,
                    "width": width,
                    "height": height,
                    "is_embedded": True
                }
                
        except Exception as e:
            self.logger.error(f"動画プロパティ取得エラー: {e}")
            return None 
    
    def _extract_videos_from_pptx_structure(self, pptx_path: str, slide_number: int) -> List[Dict[str, Any]]:
        """
        PPTXファイルをZIPとして開いてメディアファイルを直接検査
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            
        Returns:
            動画情報のリスト
        """
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            videos_info = []
            
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                # メディアファイルのリスト
                media_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/media/')]
                video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv']
                
                # 動画ファイルを検出
                video_files = [f for f in media_files if any(f.lower().endswith(ext) for ext in video_extensions)]
                
                if video_files:
                    self.logger.info(f"PPTXファイル内に {len(video_files)} 個の動画ファイルを検出")
                    
                    # スライドとメディアの関連付けを確認
                    slide_rels_path = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
                    if slide_rels_path in pptx_zip.namelist():
                        rels_content = pptx_zip.read(slide_rels_path)
                        rels_tree = ET.fromstring(rels_content)
                        
                        # 名前空間の定義
                        ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                        
                        # メディアリレーションシップを検索
                        for rel in rels_tree.findall('.//r:Relationship', ns):
                            target = rel.get('Target')
                            if target and '../media/' in target:
                                media_filename = target.split('/')[-1]
                                self.logger.info(f"スライド {slide_number} に関連付けられたメディア: {media_filename}")
                                
                                # 動画ファイルかどうかを確認
                                file_extension = media_filename.lower().split('.')[-1]
                                video_extensions = ['mp4', 'avi', 'mov', 'wmv', 'flv', 'mkv']
                                
                                if file_extension in video_extensions:
                                    # 動画情報を作成
                                    video_info = {
                                        "slide_number": slide_number,
                                        "filename": f"slide_{slide_number:02d}_video",
                                        "description": f"スライド {slide_number} の埋め込み動画",
                                        "type": "embedded_video",
                                        "is_embedded": True,
                                        "media_file": media_filename,
                                        "duration": 10,  # デフォルト値
                                        "width": 640,
                                        "height": 480
                                    }
                                    videos_info.append(video_info)
                                    self.logger.info(f"動画ファイルを検出: {media_filename}")
                                else:
                                    self.logger.debug(f"画像ファイルをスキップ: {media_filename}")
            
            return videos_info
            
        except Exception as e:
            self.logger.error(f"PPTXファイル構造検査でエラー: {e}")
            return []
    
    def _check_element_for_video(self, element) -> bool:
        """
        XML要素内に動画関連のタグがあるかチェック
        
        Args:
            element: XML要素
            
        Returns:
            動画関連のタグがあるかどうか
        """
        try:
            # より厳密な動画関連のタグ名
            video_tags = ['video', 'movie', 'av']
            
            # 要素のタグ名をチェック（より厳密に）
            element_tag = element.tag.lower()
            if any(tag in element_tag for tag in video_tags):
                # 実際に動画関連の要素かどうかをさらに確認
                if 'media' in element_tag and 'video' in element_tag:
                    return True
                elif 'movie' in element_tag:
                    return True
                elif 'av' in element_tag:
                    return True
            
            # 属性をチェック
            for attr_name, attr_value in element.attrib.items():
                attr_name_lower = attr_name.lower()
                attr_value_lower = attr_value.lower()
                
                # 動画関連の属性をチェック
                if any(tag in attr_name_lower for tag in video_tags) or any(tag in attr_value_lower for tag in video_tags):
                    return True
            
            # 子要素を再帰的にチェック（ただし、テキストボックスは除外）
            for child in element:
                child_tag = child.tag.lower()
                # テキストボックス関連の要素を除外
                if 'text' in child_tag or 'txbody' in child_tag:
                    continue
                if self._check_element_for_video(child):
                    return True
            
            return False
            
        except Exception:
            return False 
    
    def _create_video_info_from_shape(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """
        シェイプから動画情報を作成
        
        Args:
            shape: シェイプオブジェクト
            slide_number: スライド番号
            
        Returns:
            動画情報の辞書
        """
        try:
            # 基本情報
            video_info = {
                "slide_number": slide_number,
                "filename": f"slide_{slide_number:02d}_video",
                "description": f"スライド {slide_number} の埋め込み動画",
                "type": "embedded_video",
                "is_embedded": True,
                "shape_name": getattr(shape, 'name', 'unknown'),
                "shape_type": str(getattr(shape, 'shape_type', 'unknown'))
            }
            
            # サイズ情報を取得
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # EMUからピクセルに変換
                video_info['width'] = int(shape.width / 9525)  # EMU to pixels
                video_info['height'] = int(shape.height / 9525)
            else:
                video_info['width'] = 640
                video_info['height'] = 480
            
            # デフォルトの動画長
            video_info['duration'] = 10
            
            self.logger.info(f"動画情報作成: {video_info}")
            return video_info
            
        except Exception as e:
            self.logger.error(f"動画情報作成エラー: {e}")
            return None

    def _extract_embedded_video_to_mp4(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """
        埋め込み動画をMP4ファイルとして抽出（改善版）
        
        Args:
            shape: PowerPoint COMのシェイプオブジェクト
            slide_number: スライド番号
            
        Returns:
            動画情報の辞書
        """
        try:
            import win32com.client
            import tempfile
            import os
            
            # 一時ファイルのパスを生成
            temp_dir = Path(tempfile.gettempdir()) / "ppt2yt_videos"
            temp_dir.mkdir(exist_ok=True)
            
            video_filename = f"slide_{slide_number:02d}_video.mp4"
            video_path = temp_dir / video_filename
            
            # PowerPoint COMを使用して動画を抽出
            if hasattr(shape, 'MediaFormat'):
                media_format = shape.MediaFormat
                
                # 動画の長さを取得
                duration = 10  # デフォルト値
                if hasattr(media_format, 'Duration'):
                    duration = media_format.Duration / 1000  # ミリ秒から秒に変換
                
                # 動画のサイズを取得
                width = 640
                height = 480
                if hasattr(shape, 'Width') and hasattr(shape, 'Height'):
                    width = int(shape.Width)
                    height = int(shape.Height)
                
                # 方法1: 動画のソースファイルを直接取得
                try:
                    if hasattr(media_format, 'SourceFullName'):
                        source_path = media_format.SourceFullName
                        if source_path and os.path.exists(source_path):
                            # ソースファイルをコピー
                            import shutil
                            shutil.copy2(source_path, video_path)
                            self.logger.info(f"ソース動画ファイルをコピーしました: {source_path}")
                            
                            return {
                                "slide_number": slide_number,
                                "filename": video_filename,
                                "description": f"スライド {slide_number} の埋め込み動画",
                                "type": "embedded_video",
                                "is_embedded": True,
                                "video_path": str(video_path),
                                "duration": duration,
                                "width": width,
                                "height": height
                            }
                except Exception as e:
                    self.logger.debug(f"ソースファイル取得でエラー: {e}")
                
                # 方法2: 動画を一時ファイルにエクスポート（改善版）
                try:
                    # 動画を一時ファイルにエクスポート
                    export_path = str(video_path)
                    
                    # 動画専用のエクスポート方法を試行
                    if hasattr(shape, 'Export'):
                        # 動画としてエクスポート（ppMediaTypeMovie = 2）
                        shape.Export(export_path, 2)
                        
                        # ファイルが実際に作成されたかチェック
                        if video_path.exists() and video_path.stat().st_size > 0:
                            # ファイルサイズが大きい場合（動画の可能性）
                            if video_path.stat().st_size > 10000:  # 10KB以上
                                # 動画ファイルを検証
                                if self._verify_video_file(video_path):
                                    self.logger.info(f"動画ファイルを保存しました: {video_path}")
                                    
                                    return {
                                        "slide_number": slide_number,
                                        "filename": video_filename,
                                        "description": f"スライド {slide_number} の埋め込み動画",
                                        "type": "embedded_video",
                                        "is_embedded": True,
                                        "video_path": str(video_path),
                                        "duration": duration,
                                        "width": width,
                                        "height": height
                                    }
                                else:
                                    self.logger.warning(f"動画ファイルの検証に失敗（静止画の可能性）: {video_path}")
                                    return None
                            else:
                                self.logger.warning(f"動画ファイルが小さすぎます（静止画の可能性）: {video_path}")
                                return None
                        else:
                            self.logger.warning(f"動画ファイルの保存に失敗: {video_path}")
                            return None
                            
                except Exception as e:
                    self.logger.error(f"動画エクスポートでエラー: {e}")
                    return None
            
            return None
            
        except Exception as e:
            self.logger.error(f"動画抽出でエラー: {e}")
            return None

    def _verify_video_file(self, video_path: Path) -> bool:
        """
        動画ファイルが実際に動画かどうかを検証
        
        Args:
            video_path: 動画ファイルのパス
            
        Returns:
            実際の動画ファイルかどうか
        """
        try:
            import subprocess
            
            # ffprobeを使用して動画情報を取得
            cmd = [
                'ffprobe', 
                '-v', 'quiet', 
                '-print_format', 'json', 
                '-show_format', 
                '-show_streams', 
                str(video_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            
            if result.returncode == 0:
                import json
                probe_data = json.loads(result.stdout)
                
                # 動画ストリームがあるかチェック
                streams = probe_data.get('streams', [])
                video_streams = [s for s in streams if s.get('codec_type') == 'video']
                
                if video_streams:
                    # 動画の長さをチェック
                    format_info = probe_data.get('format', {})
                    duration = float(format_info.get('duration', 0))
                    
                    if duration > 1.0:  # 1秒以上の動画
                        self.logger.info(f"動画ファイルを検証しました: {video_path} (長さ: {duration:.2f}秒)")
                        return True
                    else:
                        self.logger.warning(f"動画が短すぎます: {duration:.2f}秒")
                        return False
                else:
                    self.logger.warning(f"動画ストリームが見つかりません: {video_path}")
                    return False
            else:
                self.logger.warning(f"動画ファイルの検証に失敗: {video_path}")
                return False
                
        except Exception as e:
            self.logger.error(f"動画検証でエラー: {e}")
            return False

    def _save_slide_as_video(self, pptx_path: str, slide_number: int, duration: int = 10) -> Optional[Path]:
        """
        動画を含むスライドをMP4動画として保存
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            duration: 動画の長さ（秒）
            
        Returns:
            保存された動画のパス
        """
        try:
            import win32com.client
            import pythoncom
            
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            
            try:
                # プレゼンテーションを開く
                pptx_abs_path = str(Path(pptx_path).resolve())
                presentation = powerpoint.Presentations.Open(pptx_abs_path, WithWindow=False)
                
                # 単一スライドの一時プレゼンテーションを作成
                temp_pres = powerpoint.Presentations.Add()
                
                # スライドをコピー
                presentation.Slides(slide_number).Copy()
                temp_pres.Slides.Paste()
                
                # 動画ファイル名
                video_filename = f"slide_{slide_number:02d}.mp4"
                video_path = self.output_dir / video_filename
                
                # MP4として保存（ppSaveAsMP4 = 39）
                temp_pres.SaveAs(str(video_path.resolve()), 39)
                
                temp_pres.Close()
                presentation.Close()
                
                if video_path.exists():
                    self.logger.info(f"スライド {slide_number} をMP4動画として保存: {video_path}")
                    return video_path
                else:
                    self.logger.error(f"MP4ファイルの生成に失敗: {video_path}")
                    return None
                    
            finally:
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
        except Exception as e:
            self.logger.error(f"スライドの動画保存に失敗: {e}")
            return None

    def _extract_embedded_video_properly(self, pptx_path: str, slide_number: int, shape_index: int) -> Optional[Path]:
        """
        PPTXファイル構造から直接動画ファイルを抽出
        
        Args:
            pptx_path: PPTXファイルのパス
            slide_number: スライド番号
            shape_index: シェイプのインデックス
            
        Returns:
            抽出された動画ファイルのパス
        """
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                # スライドのリレーションシップファイルを読む
                slide_rels_path = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
                
                if slide_rels_path in pptx_zip.namelist():
                    rels_content = pptx_zip.read(slide_rels_path)
                    rels_tree = ET.fromstring(rels_content)
                    
                    # 名前空間の定義
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    
                    # メディアファイルを探す
                    for rel in rels_tree.findall('.//r:Relationship', ns):
                        target = rel.get('Target')
                        if target and '../media/' in target:
                            media_path = f'ppt/media/{target.split("/")[-1]}'
                            
                            if media_path in pptx_zip.namelist():
                                # 動画ファイルかチェック
                                if any(media_path.lower().endswith(ext) for ext in ['.mp4', '.avi', '.mov', '.wmv']):
                                    # 動画を抽出
                                    video_data = pptx_zip.read(media_path)
                                    
                                    # ファイル名を生成
                                    video_filename = f"slide_{slide_number:02d}_embedded_{shape_index}.mp4"
                                    video_path = self.output_dir / video_filename
                                    
                                    # ファイルとして保存
                                    with open(video_path, 'wb') as f:
                                        f.write(video_data)
                                    
                                    self.logger.info(f"埋め込み動画を抽出: {video_path}")
                                    return video_path
            
            return None
            
        except Exception as e:
            self.logger.error(f"埋め込み動画の抽出に失敗: {e}")
            return None