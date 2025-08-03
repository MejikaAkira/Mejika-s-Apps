"""
台本生成モジュール
PPTXファイルから対話形式の台本を自動生成
"""
import json
import time
from pathlib import Path
from typing import Dict, List, Any, Optional
from pptx import Presentation
from openai import OpenAI
import re

from ..utils.config import config
from ..utils.logger import get_logger


class ScriptGenerator:
    """PPTXファイルから台本を生成するクラス"""
    
    def __init__(self):
        """ScriptGeneratorの初期化"""
        self.logger = get_logger("ScriptGenerator")
        api_key = config.get("openai.api_key")
        if not api_key:
            raise ValueError("OpenAI APIキーが設定されていません")
        # OpenAIクライアントを初期化（シンプルな初期化）
        from openai import OpenAI
        
        # 通常の初期化（プロキシ設定なし）
        self.client = OpenAI(api_key=api_key)
        self.model = config.get("openai.model.script", "gpt-4o-mini")
        self.max_tokens = config.get("openai.max_tokens", 4000)
        self.temperature = config.get("openai.temperature", 0.7)
    
    def test_connection(self) -> bool:
        """
        OpenAI APIへの接続をテスト
        
        Returns:
            接続成功フラグ
        """
        try:
            self.logger.info("OpenAI API接続をテスト中...")
            
            # 簡単なテストリクエスト
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "user", "content": "こんにちは"}
                ],
                max_tokens=10
            )
            
            self.logger.info("OpenAI API接続テスト成功")
            return True
            
        except Exception as e:
            self.logger.error(f"OpenAI API接続テスト失敗: {e}")
            return False
    
    def extract_slide_content(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        PPTXファイルからスライド内容を抽出
        
        Args:
            pptx_path: PPTXファイルのパス
            
        Returns:
            スライド内容のリスト
        """
        self.logger.info(f"PPTXファイルを読み込み中: {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "slide_number": i,
                    "text_content": [],
                    "shapes": []
                }
                
                # テキスト内容を抽出
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["text_content"].append({
                            "type": "text",
                            "content": shape.text.strip()
                        })
                    
                    # 図表や画像の情報を記録
                    if shape.shape_type in [13, 17, 19]:  # 画像、図表、SmartArt
                        slide_data["shapes"].append({
                            "type": "visual",
                            "shape_type": shape.shape_type
                        })
                
                slides_content.append(slide_data)
                self.logger.debug(f"スライド {i}: {len(slide_data['text_content'])}個のテキスト要素を抽出")
            
            self.logger.info(f"合計 {len(slides_content)} スライドの内容を抽出完了")
            return slides_content
            
        except Exception as e:
            self.logger.error(f"PPTXファイルの読み込みに失敗: {e}")
            raise
    
    def generate_dialogue_script(self, slides_content: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        対話形式の台本を生成
        
        Args:
            slides_content: スライド内容のリスト
            
        Returns:
            台本データ
        """
        self.logger.info("対話形式の台本を生成中...")
        
        # スライド内容をテキストに変換
        slides_text = self._format_slides_for_prompt(slides_content)
        
        # プロンプトを作成
        prompt = self._create_script_prompt(slides_text)
        
        try:
            # OpenAI APIで台本生成
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "あなたはプレゼンテーションの台本作成の専門家です。"},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=self.max_tokens,
                temperature=self.temperature
            )
            
            script_text = response.choices[0].message.content
            script_data = self._parse_script_response(script_text, slides_content)
            
            self.logger.info("台本生成完了")
            return script_data
            
        except Exception as e:
            self.logger.error(f"OpenAI APIでの台本生成に失敗: {e}")
            self.logger.warning("フォールバック台本を生成します")
            
            # フォールバック台本を生成
            return self._create_fallback_script(slides_content)
    
    def _format_slides_for_prompt(self, slides_content: List[Dict[str, Any]]) -> str:
        """
        スライド内容をプロンプト用にフォーマット
        
        Args:
            slides_content: スライド内容のリスト
            
        Returns:
            フォーマットされたテキスト
        """
        formatted_text = "プレゼンテーション内容:\n\n"
        
        for slide in slides_content:
            formatted_text += f"【スライド {slide['slide_number']}】\n"
            
            for text_item in slide['text_content']:
                formatted_text += f"- {text_item['content']}\n"
            
            if slide['shapes']:
                formatted_text += f"- 図表・画像: {len(slide['shapes'])}個\n"
            
            formatted_text += "\n"
        
        return formatted_text
    
    def _create_script_prompt(self, slides_text: str) -> str:
        """
        台本生成用のプロンプトを作成
        
        Args:
            slides_text: スライド内容のテキスト
            
        Returns:
            プロンプトテキスト
        """
        return f"""
以下のプレゼンテーション内容を基に、自然な対話形式の台本を作成してください。

{slides_text}

重要要件:
1. **全スライドを必ずカバーしてください** - スライド1から最後のスライドまで、すべてのスライドにセリフを割り当ててください
2. メインスピーカー（男性）とアシスタント（女性）の2人による自然な対話形式
3. メインスピーカーが内容を説明し、アシスタントが質問や相槌、理解を示す反応をする
4. 各セリフの長さに応じて適切な表示時間を設定（1文字約0.1秒、最小2秒、最大15秒）
5. 各スライドの内容を適切に説明
6. 専門用語は分かりやすく説明
7. アシスタントのキャラクター設定：
   - 天然ボケだけど鋭い洞察力を持ち、専門用語を分かりやすい言葉に咀嚼してくれる
   - 明るくエネルギッシュな性格で、優しくて親しみやすい性格
   - 以下のような反応をする：
     - 相槌：「なるほど〜そうなんですね〜」「へえ〜それは知らなかったです」
     - 質問：「それって、つまりどういうことですか？」「もっと簡単に言うと？」
     - 理解：「あ、なるほど！理解しました〜」「分かりました〜」
     - 興味：「面白いですね〜」「もっと詳しく教えてください〜」
     - 専門用語の咀嚼：「つまり、○○ということですね？」
8. 各スライドの説明が終わったら、次のスライドに自然に移行
9. スライドの切り替えは内容の区切りに合わせて行う
10. メインスピーカーとアシスタントの役割は固定（入れ替わらない）

スライドカバー要件:
- スライド1から最後のスライドまで、すべてのスライドに最低1つ以上のセリフを割り当ててください
- 内容が少ないスライドでも、メインスピーカーが簡単に説明し、アシスタントが相槌や質問をする形でカバーしてください
- スライドの内容が空でも、スライド番号を確認して適切にセリフを割り当ててください
- 図がある場合、図の説明を含めてください

対話の例:
- メインスピーカー: 「RMSについて説明します」
- アシスタント: 「RMSって何ですか？もっと簡単に言うと？」
- メインスピーカー: 「RMSは...」
- アシスタント: 「あ、なるほど！つまり、音の大きさを表す数字ということですね〜」

出力形式:
```json
{{
  "title": "プレゼンテーションタイトル",
  "total_duration": 300,
  "dialogue": [
    {{
      "slide_file": "slide_01.png",
      "timestamp": "00:00:00",
      "text": "セリフ内容",
      "duration": 8,
      "role": "main_speaker"
    }}
  ]
}}
```

注意:
- **全スライドを必ずカバーしてください** - スライド1から最後まで、すべてのスライドにセリフを割り当ててください
- roleフィールドで「main_speaker」（メインスピーカー）または「assistant」（アシスタント）を指定してください
- slide_fileフィールドで対応する画像ファイル名を指定してください（slide_01.png, slide_02.png, ...）
- durationは実際のセリフの長さに応じて計算してください
- 短い相槌は2-3秒、質問は3-5秒、説明は5-15秒程度
- メインスピーカーが主に説明し、アシスタントが質問や相槌をする自然な流れにしてください
- アシスタントは質問や相槌、理解を示す自然な反応をしてください
- 各スライドの説明が完了したら、次のスライドに移行してください
- システムが各スライドの説明終了後に自然な間を自動で追加します
"""
    
    def _parse_script_response(self, script_text: str, slides_content: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        生成された台本テキストをパース
        
        Args:
            script_text: 生成された台本テキスト
            slides_content: 元のスライド内容
            
        Returns:
            パースされた台本データ
        """
        try:
            # JSON部分を抽出
            json_start = script_text.find('{')
            json_end = script_text.rfind('}') + 1
            
            if json_start == -1 or json_end == 0:
                raise ValueError("JSON形式の台本が見つかりません")
            
            json_text = script_text[json_start:json_end]
            script_data = json.loads(json_text)
            
            # speakerフィールドを自動割り当て（交互に）
            self._assign_speakers(script_data)
            
            # タイムスタンプを計算
            self._calculate_timestamps(script_data)
            
            # 全スライドカバー検証
            self._validate_slide_coverage(script_data, slides_content)
            
            return script_data
            
        except Exception as e:
            self.logger.error(f"台本のパースに失敗: {e}")
            # フォールバック台本を作成
            return self._create_fallback_script(slides_content)
    
    def _assign_speakers(self, script_data: Dict[str, Any]):
        """
        台本の各セリフにspeakerとvoiceをroleに基づいて割り当て
        
        Args:
            script_data: 台本データ
        """
        # 設定ファイルから音声設定を読み込み
        male_voice = config.get("openai.tts.male_voice", "echo")
        female_voice = config.get("openai.tts.female_voice", "alloy")
        
        dialogue = script_data.get("dialogue", [])
        
        for i, line in enumerate(dialogue):
            role = line.get("role", "main_speaker")  # デフォルトはメインスピーカー
            
            if role == "main_speaker":
                line["speaker"] = "speaker"
                line["voice"] = male_voice
                line["voice_type"] = "male"
            elif role == "assistant":
                line["speaker"] = "listener"
                line["voice"] = female_voice
                line["voice_type"] = "female"
            else:
                # 不明な役割の場合はメインスピーカーとして扱う
                line["speaker"] = "speaker"
                line["voice"] = male_voice
                line["voice_type"] = "male"
                self.logger.warning(f"不明な役割 '{role}' をメインスピーカーとして扱います")
            
            self.logger.debug(f"セリフ {i+1}: {line['speaker']} ({line['voice']}) - {line['text'][:30]}...")
    
    def _calculate_timestamps(self, script_data: Dict[str, Any]):
        """
        台本のタイムスタンプを計算
        
        Args:
            script_data: 台本データ
        """
        current_time = 0
        
        for dialogue in script_data.get("dialogue", []):
            dialogue["timestamp"] = self._seconds_to_timestamp(current_time)
            current_time += dialogue.get("duration", 5)
        
        script_data["total_duration"] = current_time
    
    def _seconds_to_timestamp(self, seconds: int) -> str:
        """
        秒数をタイムスタンプ形式に変換
        
        Args:
            seconds: 秒数
            
        Returns:
            タイムスタンプ文字列 (HH:MM:SS)
        """
        hours = seconds // 3600
        minutes = (seconds % 3600) // 60
        secs = seconds % 60
        return f"{hours:02d}:{minutes:02d}:{secs:02d}"
    
    def _create_fallback_script(self, slides_content: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        フォールバック用の基本的な台本を作成
        
        Args:
            slides_content: スライド内容
            
        Returns:
            基本的な台本データ
        """
        self.logger.warning("フォールバック台本を作成します")
        
        dialogue = []
        current_time = 0
        
        for slide in slides_content:
            # メインスピーカーの説明
            slide_text = " ".join([item["content"] for item in slide["text_content"]])
            
            dialogue.append({
                "slide_file": f"slide_{slide['slide_number']:02d}.png",
                "timestamp": self._seconds_to_timestamp(current_time),
                "text": f"スライド{slide['slide_number']}について説明します。{slide_text}",
                "duration": 10,
                "role": "main_speaker"
            })
            current_time += 10
            
            # アシスタントの反応（天然ボケキャラ）
            reactions = [
                "あ、なるほど〜！理解しました〜",
                "そうなんですね〜、分かりました〜",
                "へえ〜、面白いですね〜",
                "もっと詳しく教えてください〜",
                "それって、つまりどういうことですか？"
            ]
            
            dialogue.append({
                "slide_file": f"slide_{slide['slide_number']:02d}.png",
                "timestamp": self._seconds_to_timestamp(current_time),
                "text": reactions[slide["slide_number"] % len(reactions)],
                "duration": 3,
                "role": "assistant"
            })
            current_time += 3
        
        script_data = {
            "title": "プレゼンテーション",
            "total_duration": current_time,
            "dialogue": dialogue
        }
        
        # speakerとvoiceを自動割り当て
        self._assign_speakers(script_data)
        
        return script_data
    
    def save_script(self, script_data: Dict[str, Any], output_path: str):
        """
        台本をファイルに保存
        
        Args:
            script_data: 台本データ
            output_path: 出力ファイルパス
        """
        try:
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(script_data, f, ensure_ascii=False, indent=2)
            
            self.logger.info(f"台本を保存しました: {output_path}")
            
        except Exception as e:
            self.logger.error(f"台本の保存に失敗: {e}")
            raise
    
    def generate_script(self, pptx_path: str, output_path: str, videos_info: Optional[List[Dict[str, Any]]] = None) -> Dict[str, Any]:
        """
        PPTXファイルから台本を生成
        
        Args:
            pptx_path: PPTXファイルのパス
            output_path: 出力ファイルのパス
            videos_info: 動画情報のリスト（オプション）
            
        Returns:
            生成された台本データ
        """
        try:
            # スライド内容を抽出
            slides_content = self.extract_slide_content(pptx_path)
            
            # 台本を生成
            script_data = self.generate_dialogue_script(slides_content)
            
            # 動画情報がある場合は台本を調整
            if videos_info and config.get("video_processing.adjust_script_for_videos", True):
                script_data = self._adjust_script_for_videos(script_data, videos_info)
            
            # 台本を保存
            self.save_script(script_data, output_path)
            
            # 生成結果のログ出力
            self._log_script_summary(script_data)
            
            return script_data
            
        except Exception as e:
            self.logger.error(f"台本生成処理でエラー: {e}")
            raise
    
    def _log_script_summary(self, script_data: Dict[str, Any]):
        """
        生成された台本のサマリーをログ出力
        
        Args:
            script_data: 台本データ
        """
        dialogue = script_data.get("dialogue", [])
        male_lines = sum(1 for line in dialogue if line.get("voice_type") == "male")
        female_lines = sum(1 for line in dialogue if line.get("voice_type") == "female")
        
        # 音声設定を取得
        male_voice = config.get("openai.tts.male_voice", "echo")
        female_voice = config.get("openai.tts.female_voice", "alloy")
        
        self.logger.info(f"台本生成完了:")
        self.logger.info(f"  タイトル: {script_data.get('title', 'N/A')}")
        self.logger.info(f"  総再生時間: {script_data.get('total_duration', 0)}秒")
        self.logger.info(f"  総セリフ数: {len(dialogue)}")
        self.logger.info(f"  男性セリフ: {male_lines} ({male_voice})")
        self.logger.info(f"  女性セリフ: {female_lines} ({female_voice})")
        
        # 最初の数行をサンプル表示
        self.logger.info("サンプルセリフ:")
        for i, line in enumerate(dialogue[:4]):
            self.logger.info(f"  {i+1}. [{line.get('speaker', 'N/A')}] {line.get('text', 'N/A')[:50]}...") 

    def _validate_slide_coverage(self, script_data: Dict[str, Any], slides_content: List[Dict[str, Any]]):
        """
        全スライドがカバーされているかを検証し、不足があれば補完
        
        Args:
            script_data: 台本データ
            slides_content: 元のスライド内容
        """
        try:
            # カバーされているスライド番号を抽出
            covered_slides = set()
            for dialogue in script_data.get('dialogue', []):
                slide_file = dialogue.get('slide_file', '')
                match = re.search(r'slide_(\d+)\.png', slide_file)
                if match:
                    slide_num = int(match.group(1))
                    covered_slides.add(slide_num)
            
            # 全スライド番号を取得
            all_slides = set(slide['slide_number'] for slide in slides_content)
            
            # 不足しているスライドを特定
            missing_slides = all_slides - covered_slides
            
            if missing_slides:
                self.logger.warning(f"以下のスライドがカバーされていません: {sorted(missing_slides)}")
                self.logger.info("不足しているスライドのセリフを補完します")
                
                # 不足しているスライドのセリフを補完
                self._add_missing_slide_dialogues(script_data, missing_slides, slides_content)
                
                # タイムスタンプを再計算
                self._calculate_timestamps(script_data)
                
                self.logger.info("スライドカバー補完完了")
            else:
                self.logger.info("全スライドがカバーされています")
                
        except Exception as e:
            self.logger.error(f"スライドカバレッジ検証でエラー: {e}")

    def _adjust_script_for_videos(self, script_data: Dict[str, Any], videos_info: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        動画情報に基づいて台本を調整
        
        Args:
            script_data: 台本データ
            videos_info: 動画情報のリスト
            
        Returns:
            調整された台本データ
        """
        try:
            self.logger.info("動画情報に基づいて台本を調整中...")
            
            # 動画情報をスライド番号でグループ化
            videos_by_slide = {}
            for video_info in videos_info:
                slide_number = video_info.get('slide_number', 0)
                if slide_number not in videos_by_slide:
                    videos_by_slide[slide_number] = []
                videos_by_slide[slide_number].append(video_info)
            
            # 台本の各セリフを調整
            adjusted_dialogue = []
            for dialogue in script_data.get('dialogue', []):
                slide_file = dialogue.get('slide_file', '')
                match = re.search(r'slide_(\d+)\.png', slide_file)
                
                if match:
                    slide_number = int(match.group(1))
                    
                    # このスライドに動画があるかチェック
                    if slide_number in videos_by_slide:
                        # 動画の長さを取得
                        video_duration = self._get_total_video_duration(videos_by_slide[slide_number])
                        
                        # 動画の長さに合わせてセリフを調整
                        adjusted_dialogue.extend(self._create_video_adjusted_dialogues(dialogue, video_duration))
                        self.logger.info(f"スライド {slide_number} の動画長に合わせて調整: {video_duration}秒")
                    else:
                        # 動画がない場合は通常のセリフ
                        adjusted_dialogue.append(dialogue)
                else:
                    # スライド番号が特定できない場合は通常のセリフ
                    adjusted_dialogue.append(dialogue)
            
            # 調整された台本を更新
            script_data['dialogue'] = adjusted_dialogue
            
            # タイムスタンプを再計算
            self._calculate_timestamps(script_data)
            
            self.logger.info(f"動画調整完了: {len(adjusted_dialogue)}セリフ")
            return script_data
            
        except Exception as e:
            self.logger.error(f"動画調整でエラー: {e}")
            return script_data
    
    def _get_total_video_duration(self, videos: List[Dict[str, Any]]) -> float:
        """
        複数の動画の合計長を取得
        
        Args:
            videos: 動画情報のリスト
            
        Returns:
            合計動画長（秒）
        """
        total_duration = 0.0
        for video in videos:
            duration = video.get('duration', 0)
            total_duration += duration
        return total_duration
    
    def _create_video_adjusted_dialogues(self, original_dialogue: Dict[str, Any], video_duration: float) -> List[Dict[str, Any]]:
        """
        動画の長さに合わせてセリフを調整
        
        Args:
            original_dialogue: 元のセリフ
            video_duration: 動画の長さ（秒）
            
        Returns:
            調整されたセリフのリスト
        """
        # 設定から閾値を取得
        short_threshold = config.get("video_processing.video_script_adjustment.short_video_threshold", 10)
        long_threshold = config.get("video_processing.video_script_adjustment.long_video_threshold", 15)
        min_duration = config.get("video_processing.video_script_adjustment.min_dialogue_duration", 3)
        
        adjusted_dialogues = []
        
        # 動画の長さに基づいてセリフを分割または調整
        if video_duration > short_threshold:  # 長い動画の場合
            # 動画の開始、中間、終了でセリフを分割
            original_text = original_dialogue.get('text', '')
            
            # 開始セリフ
            start_dialogue = original_dialogue.copy()
            start_dialogue['text'] = f"この動画をご覧ください。{original_text}"
            start_dialogue['duration'] = min(min_duration, video_duration * 0.2)
            adjusted_dialogues.append(start_dialogue)
            
            # 中間セリフ（動画が長い場合のみ）
            if video_duration > long_threshold:
                middle_dialogue = original_dialogue.copy()
                middle_dialogue['text'] = "動画が続いています。"
                middle_dialogue['duration'] = min(2.0, video_duration * 0.1)
                adjusted_dialogues.append(middle_dialogue)
            
            # 終了セリフ
            end_dialogue = original_dialogue.copy()
            end_dialogue['text'] = f"動画が終了しました。{original_text}"
            end_dialogue['duration'] = min(min_duration, video_duration * 0.2)
            adjusted_dialogues.append(end_dialogue)
            
        else:  # 短い動画の場合
            # 元のセリフを動画の長さに合わせて調整
            adjusted_dialogue = original_dialogue.copy()
            adjusted_dialogue['duration'] = max(original_dialogue.get('duration', 5), video_duration + 2)
            adjusted_dialogue['text'] = f"この動画と合わせてご覧ください。{original_dialogue.get('text', '')}"
            adjusted_dialogues.append(adjusted_dialogue)
        
        return adjusted_dialogues 
    
    def _add_missing_slide_dialogues(self, script_data: Dict[str, Any], missing_slides: set, slides_content: List[Dict[str, Any]]):
        """
        不足しているスライドのセリフを補完
        
        Args:
            script_data: 台本データ
            missing_slides: 不足しているスライド番号のセット
            slides_content: 元のスライド内容
        """
        dialogue = script_data.get('dialogue', [])
        
        for slide_num in sorted(missing_slides):
            # 該当するスライドの内容を取得
            slide_content = next((slide for slide in slides_content if slide['slide_number'] == slide_num), None)
            
            if slide_content:
                # スライドの内容をテキストに変換
                slide_text = " ".join([item["content"] for item in slide_content["text_content"]])
                
                # メインスピーカーの説明セリフを追加
                dialogue.append({
                    "slide_file": f"slide_{slide_num:02d}.png",
                    "timestamp": "00:00:00",  # 後で再計算される
                    "text": f"スライド{slide_num}について説明します。{slide_text}",
                    "duration": 8,
                    "role": "main_speaker",
                    "speaker": "speaker",
                    "voice": "onyx",
                    "voice_type": "male"
                })
                
                # アシスタントの反応セリフを追加
                reactions = [
                    "なるほど〜、理解しました〜",
                    "そうなんですね〜、分かりました〜",
                    "へえ〜、面白いですね〜",
                    "もっと詳しく教えてください〜",
                    "それって、つまりどういうことですか？"
                ]
                
                dialogue.append({
                    "slide_file": f"slide_{slide_num:02d}.png",
                    "timestamp": "00:00:00",  # 後で再計算される
                    "text": reactions[slide_num % len(reactions)],
                    "duration": 3,
                    "role": "assistant",
                    "speaker": "listener",
                    "voice": "alloy",
                    "voice_type": "female"
                })
                
                self.logger.info(f"スライド{slide_num}のセリフを補完しました")
        
        script_data['dialogue'] = dialogue 